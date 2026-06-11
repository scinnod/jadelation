# SPDX-License-Identifier: Apache-2.0
# SPDX-FileCopyrightText: 2024-2026 David Kleinhans, Jade University of Applied Sciences
"""
Views for DeepL translation frontend application.

This module provides views for the translation interface and usage statistics.
"""

import logging
import os
import re
import threading
import time
from datetime import timedelta
from dateutil.relativedelta import relativedelta

import deepl
from django.conf import settings
from django.db.models import Count, F, Sum
from django.http import HttpResponse, FileResponse, JsonResponse
from django.shortcuts import render
from django.utils import timezone
from django.utils.translation import gettext as _, get_language, activate
from django.views.decorators.http import require_POST, require_GET

from .forms import TranslationForm, DocumentTranslationForm
from .models import Translation, Glossary, DocumentTranslationJob

# Configure logging
logger = logging.getLogger(__name__)


# Load configuration constants from settings (with defaults as fallback)
TRANSLATION_DETECTION_LEN = getattr(settings, 'TRANSLATION_DETECTION_LEN', 30)
STATISTICS_DAYS = getattr(settings, 'STATISTICS_DAYS', 31)
STATISTICS_MONTHS = getattr(settings, 'STATISTICS_MONTHS', 24)
STATISTICS_YEARS = getattr(settings, 'STATISTICS_YEARS', 5)

# DeepL model type for auto-detection vs. actual translation.
# Empty string means "let the API decide" (omit the parameter).
DEEPL_MODEL_TYPE_DETECTION = getattr(settings, 'DEEPL_MODEL_TYPE_DETECTION', '') or None
DEEPL_MODEL_TYPE_TRANSLATION = getattr(settings, 'DEEPL_MODEL_TYPE_TRANSLATION', '') or None

# Time-to-live for the glossary cache in seconds (default: 1 hour).
# The cache is refreshed lazily on the next translation request after expiry.
GLOSSARY_CACHE_TTL = getattr(settings, 'GLOSSARY_CACHE_TTL', 3600)

# ---------------------------------------------------------------------------
# Document translation feature helpers
# ---------------------------------------------------------------------------

def _is_doc_translation_possible():
    """Return ``True`` when document translation *might* be used by someone.

    Used for infrastructure tasks (cleanup) that should run whenever the
    feature is ``True`` or a regex string (i.e. not entirely disabled).
    """
    val = getattr(settings, 'DOCUMENT_TRANSLATION_ENABLED', False)
    return val is not False


def _is_doc_translation_enabled(request):
    """Return ``True`` when document translation is enabled for *request*.

    Resolves the tri-state ``DOCUMENT_TRANSLATION_ENABLED`` setting:
    - ``False``  → always disabled
    - ``True``   → always enabled
    - str (regex)→ enabled only when the middleware set the session flag
    """
    val = getattr(settings, 'DOCUMENT_TRANSLATION_ENABLED', False)
    if val is False:
        return False
    if val is True:
        return True
    return bool(request.session.get('document_translations', False))


class _GlossaryCache:
    """
    Lazy, auto-refreshing glossary cache.

    Instead of loading glossaries once at module-import time (which
    silently fails when the DeepL API or the database is not yet
    available during Docker startup), the cache is populated on the
    first translation request and refreshed after *GLOSSARY_CACHE_TTL*
    seconds.

    The cache maps normalised 2-letter language pairs (e.g. ``"DE->EN"``)
    to ``deepl.GlossaryInfo`` objects fetched from the DeepL API.
    """

    def __init__(self):
        self._cache = {}
        self._loaded_at = float('-inf')  # guarantees first _ensure_loaded() triggers

    # ------------------------------------------------------------------
    # public helpers
    # ------------------------------------------------------------------

    def get(self, key, default=None):
        """Look up a glossary by normalised language-pair key."""
        self._ensure_loaded()
        return self._cache.get(key, default)

    def reload(self):
        """Force an immediate reload of the cache."""
        self._load()

    @property
    def loaded_pairs(self):
        """Return the set of cached language-pair keys (for diagnostics)."""
        self._ensure_loaded()
        return set(self._cache.keys())

    # ------------------------------------------------------------------
    # internals
    # ------------------------------------------------------------------

    def _ensure_loaded(self):
        """Load (or refresh) the cache when stale or empty."""
        if time.monotonic() - self._loaded_at > GLOSSARY_CACHE_TTL:
            self._load()

    def _load(self):
        """
        Fetch GlossaryInfo objects from the DeepL API for every glossary
        stored in the local database.

        Use management commands to manage glossaries:
         - ``python manage.py glossary_put <csv> <name> <source> <target>``
         - ``python manage.py glossary_list``
         - ``python manage.py glossary_remove <name_or_id>``
        """
        new_cache = {}

        try:
            translator = deepl.Translator(settings.DEEPL_AUTHKEY)
            db_glossaries = Glossary.objects.all()

            for db_glossary in db_glossaries:
                try:
                    deepl_glossary = translator.get_glossary(
                        db_glossary.glossary_id
                    )

                    # Normalise to 2-letter base codes so that lookups
                    # from the translation view always match, regardless
                    # of whether the glossary was uploaded with a regional
                    # variant code (e.g. EN-GB) or a plain code (EN).
                    norm_source = db_glossary.source_lang[:2].upper()
                    norm_target = db_glossary.target_lang[:2].upper()
                    cache_key = f"{norm_source}->{norm_target}"

                    # The queryset is ordered by -upload_date (newest
                    # first).  Keep only the most recent glossary per
                    # language pair and skip older duplicates.
                    if cache_key in new_cache:
                        logger.info(
                            "Skipping older duplicate glossary: %s (%s)",
                            db_glossary.name,
                            cache_key,
                        )
                        continue

                    new_cache[cache_key] = deepl_glossary

                    logger.info(
                        "Loaded glossary: %s (%s, %d entries, ready=%s)",
                        db_glossary.name,
                        cache_key,
                        db_glossary.entry_count,
                        deepl_glossary.ready,
                    )
                except deepl.DeepLException as e:
                    logger.warning(
                        "Could not load glossary %s (ID: %s): %s",
                        db_glossary.name,
                        db_glossary.glossary_id,
                        e,
                    )
                except Exception as e:
                    logger.exception(
                        "Unexpected error loading glossary %s: %s",
                        db_glossary.name,
                        e,
                    )

        except Exception as e:
            logger.error("Failed to initialise glossary cache: %s", e)

        self._cache = new_cache
        self._loaded_at = time.monotonic()
        logger.info(
            "Glossary cache refreshed: %d pair(s) loaded – %s",
            len(new_cache),
            list(new_cache.keys()) or "(none)",
        )


# Singleton instance – populated lazily on first translation request
glossary_cache = _GlossaryCache()


def _active_job_context(request):
    """Return template context dict for any active document-translation job.

    Checks whether the current session has a PENDING, PROCESSING, or recent
    COMPLETED/FAILED job and returns a JSON-safe dict so the frontend can
    resume polling or show the result immediately on page load.
    """
    if not _is_doc_translation_enabled(request):
        return {}
    session_key = getattr(request.session, 'session_key', None)
    if not session_key:
        return {}
    # Look for an active (unfinished) job first, then fall back to most-recent
    # finished job that still has a result_path (i.e. not cleaned up yet).
    job = (
        DocumentTranslationJob.objects.filter(
            session_key=session_key,
            status__in=[
                DocumentTranslationJob.Status.PENDING,
                DocumentTranslationJob.Status.PROCESSING,
            ],
        ).first()
    )
    if job is None:
        # Check for a recent COMPLETED job that has not been downloaded
        # or cleaned up yet (so the user can still access the result).
        job = (
            DocumentTranslationJob.objects.filter(
                session_key=session_key,
                status=DocumentTranslationJob.Status.COMPLETED,
                downloaded=False,
            )
            .exclude(result_path="")
            .order_by("-completed_at")
            .first()
        )
    if job is None:
        return {}
    import json
    data = {
        "id": str(job.id),
        "status": job.status,
        "original_filename": job.original_filename,
    }
    if job.status == DocumentTranslationJob.Status.COMPLETED:
        data.update({
            "characters": job.characters,
            "duration_seconds": job.duration_seconds,
            "file_size": job.file_size,
        })
    elif job.status == DocumentTranslationJob.Status.FAILED:
        data["error_message"] = job.error_message or ""
    return {"active_doc_job_json": json.dumps(data)}


def deepl_translation(request):
    """
    Handle DeepL translation requests.
    
    Supports automatic language detection and translation between German and English.
    Logs all translations to the database for usage tracking.
    """
    # Lazily clean up expired document translation files
    if _is_doc_translation_possible():
        try:
            _cleanup_stale_document_jobs()
        except Exception:
            pass  # never let cleanup break the main page

    if request.method == "POST":
        form = TranslationForm(request.POST)
        if form.is_valid():
            # Initial values of variables
            character_count = 0
            auto_detection = False
            
            # Initialize translator (catch configuration errors)
            try:
                translator = deepl.Translator(settings.DEEPL_AUTHKEY)
            except (AttributeError, ValueError) as e:
                logger.error(f"DeepL translator initialization failed: {e}")
                return HttpResponse(
                    _("ERROR: Translation service configuration error. Please contact the administrator.")
                )
            
            direction = form.cleaned_data["directionChoice"]

            # Perform auto detection if direction is "auto"
            if direction == "auto":
                try:
                    result = translator.translate_text(
                        form.cleaned_data["sourceText"][:TRANSLATION_DETECTION_LEN],
                        target_lang="DE",
                        model_type=DEEPL_MODEL_TYPE_DETECTION,
                    )
                    character_count += min(
                        len(form.cleaned_data["sourceText"]), 
                        TRANSLATION_DETECTION_LEN
                    )
                    auto_detection = True
                    
                    if result.detected_source_lang[:2] == "DE":
                        direction = "DE->EN-GB|"
                    elif result.detected_source_lang[:2] == "EN":
                        direction = "EN-GB->DE|more"
                        
                except deepl.DeepLException as e:
                    logger.error(f"DeepL API error during auto-detection: {e}")
                    return HttpResponse(
                        _(
                            "ERROR: The API services at DeepL are currently unavailable. "
                            "Please try again later."
                        )
                    )
                except Exception as e:
                    logger.exception(f"Unexpected error during auto-detection: {e}")
                    return HttpResponse(
                        _("ERROR: An unexpected error occurred. Please try again later.")
                    )

            # Perform translation if direction known
            if direction != "auto":
                
                # Parse direction string
                direction_parts = direction.split("|")
                lang_pair = direction_parts[0].split("->")
                source_lang = lang_pair[0][:2]  # Only first two letters
                target_lang = lang_pair[1]
                formality = direction_parts[1] if len(direction_parts[1]) > 0 else None
                
                # Look up glossary using normalised 2-letter codes
                # (e.g. "DE->EN" not "DE->EN-GB") to match the cache keys.
                # The DeepL API requires source_lang when a glossary is used.
                glossary_key = f"{source_lang}->{target_lang[:2]}"
                active_glossary = glossary_cache.get(glossary_key)

                if active_glossary is not None:
                    logger.info(
                        "Using glossary '%s' (id=%s) for %s",
                        active_glossary.name,
                        active_glossary.glossary_id,
                        glossary_key,
                    )
                else:
                    logger.debug(
                        "No glossary found for %s (available: %s)",
                        glossary_key,
                        glossary_cache.loaded_pairs,
                    )

                try:
                    result = translator.translate_text(
                        form.cleaned_data["sourceText"],
                        source_lang=source_lang,
                        target_lang=target_lang,
                        formality=formality,
                        glossary=active_glossary,
                        model_type=DEEPL_MODEL_TYPE_TRANSLATION,
                    )
                    translation = result.text
                    character_count += len(form.cleaned_data["sourceText"])
                    
                except deepl.DeepLException as e:
                    logger.error(f"DeepL API error during translation: {e}")
                    # Record detection characters that were already consumed
                    if character_count > 0:
                        try:
                            Translation.objects.create(
                                characters=character_count,
                                direction=direction,
                                auto_detection=auto_detection,
                            )
                        except Exception:
                            pass
                    return HttpResponse(
                        _(
                            "ERROR: The API services at DeepL are currently unavailable. "
                            "Please try again later."
                        )
                    )
                except Exception as e:
                    logger.exception(f"Unexpected error during translation: {e}")
                    if character_count > 0:
                        try:
                            Translation.objects.create(
                                characters=character_count,
                                direction=direction,
                                auto_detection=auto_detection,
                            )
                        except Exception:
                            pass
                    return HttpResponse(
                        _("ERROR: An unexpected error occurred. Please try again later.")
                    )
            else:
                translation = _(
                    "ERROR: Language recognition was not successful. "
                    "Please provide a longer excerpt or specify the direction explicitly."
                )

            # Prepare form to return (including update of direction)
            return_form = TranslationForm(
                initial={
                    "sourceText": form.cleaned_data["sourceText"],
                    "directionChoice": direction,
                    "translatedText": translation,
                }
            )

            # Log translation metadata to database
            try:
                Translation.objects.create(
                    characters=character_count,
                    direction=direction,
                    auto_detection=auto_detection,
                )
            except Exception as e:
                logger.exception(f"Failed to log translation to database: {e}")
                # Continue - don't fail the request if logging fails

            # Request current API usage and limit.
            # Note: get_usage() calls /v2/usage which returns the
            # **account-level** character usage and subscription limit,
            # not a per-API-key limit.
            try:
                usage = translator.get_usage()
                usage_str = _(
                    "Current total API character usage: {count} of {limit}."
                ).format(
                    count=usage.character.count, 
                    limit=usage.character.limit
                )
            except deepl.DeepLException as e:
                logger.warning(f"Could not retrieve DeepL usage statistics: {e}")
                usage_str = _(
                    "Current total API character usage: not available due to connection issues."
                )
            except Exception as e:
                logger.exception(f"Unexpected error retrieving usage stats: {e}")
                usage_str = _("Current total API character usage: not available.")

            return render(
                request,
                "translation_frontend.html",
                {"form": return_form, "usage": usage_str, "doc_form": DocumentTranslationForm(),
                 **_active_job_context(request)},
            )
        else:
            # Form is not valid - render with errors
            logger.warning(f"Invalid form submission: {form.errors}")

    # GET request or other method - create a blank form
    else:
        form = TranslationForm(initial={"directionChoice": "auto"})

    return render(
        request,
        "translation_frontend.html",
        {"form": form, "doc_form": DocumentTranslationForm(),
         **_active_job_context(request)},
    )


def deepl_daily_statistics(request, granularity):
    """
    Display usage statistics with different granularities (daily, monthly, yearly).
    
    Args:
        request: HTTP request object
        granularity: 'd' for daily, 'm' for monthly, 'y' for yearly
    
    Returns:
        Rendered statistics page with usage data
    """
    
    def add_line(statistics, fdt, tdt):
        """Helper function to add a statistics line for a time period."""
        req = Translation.objects.filter(
            timestamp__gte=fdt, 
            timestamp__lt=tdt
        ).aggregate(cumul=Sum("characters"), events=Count("characters"))

        req_ende = Translation.objects.filter(
            direction__startswith="EN", 
            timestamp__gte=fdt, 
            timestamp__lt=tdt
        ).aggregate(cumul=Sum("characters"), events=Count("characters"))

        req_deen = Translation.objects.filter(
            direction__startswith="DE", 
            timestamp__gte=fdt, 
            timestamp__lt=tdt
        ).aggregate(cumul=Sum("characters"), events=Count("characters"))

        statistics.append(
            {
                "fdt": fdt,
                "tdt": tdt,
                "total": "{} ({})".format(
                    req["events"], 
                    req["cumul"] if req["cumul"] else 0
                ),
                "deen": "{} ({})".format(
                    req_deen["events"], 
                    req_deen["cumul"] if req_deen["cumul"] else 0
                ),
                "ende": "{} ({})".format(
                    req_ende["events"], 
                    req_ende["cumul"] if req_ende["cumul"] else 0
                ),
            }
        )

    statistics = []
    now = timezone.now()
    
    # Generate statistics based on granularity
    if granularity == "d":
        granularity_title = _("Report of daily usage:")
        dformat = "Y-m-d"
        for days in range(STATISTICS_DAYS):
            fdt = (now - relativedelta(days=days)).replace(
                minute=0, hour=0, second=0, microsecond=0
            )
            tdt = (now - relativedelta(days=(days - 1))).replace(
                minute=0, hour=0, second=0, microsecond=0
            )
            add_line(statistics, fdt, tdt)
            
    elif granularity == "m":
        granularity_title = _("Report of monthly usage:")
        dformat = "Y-m"
        for months in range(STATISTICS_MONTHS):
            fdt = (now - relativedelta(months=months)).replace(
                day=1, minute=0, hour=0, second=0, microsecond=0
            )
            tdt = (now - relativedelta(months=(months - 1))).replace(
                day=1, minute=0, hour=0, second=0, microsecond=0
            )
            add_line(statistics, fdt, tdt)
            
    elif granularity == "y":
        granularity_title = _("Report of annual usage:")
        dformat = "Y"
        for years in range(STATISTICS_YEARS):
            fdt = (now - relativedelta(years=years)).replace(
                month=1, day=1, minute=0, hour=0, second=0, microsecond=0
            )
            tdt = (now - relativedelta(years=(years - 1))).replace(
                month=1, day=1, minute=0, hour=0, second=0, microsecond=0
            )
            add_line(statistics, fdt, tdt)
    else:
        logger.warning(f"Invalid granularity requested: {granularity}")
        granularity_title = _("Requested report not available.")
        dformat = "Y-m-d"

    # Retrieve usage statistics from DeepL API.
    # Note: get_usage() returns the account-level character usage and
    # subscription limit, not a per-API-key limit.
    try:
        translator = deepl.Translator(settings.DEEPL_AUTHKEY)
        usage = translator.get_usage()
        usage_str = _(
            "Current total API character usage: {count} of {limit}."
        ).format(
            count=usage.character.count, 
            limit=usage.character.limit
        )
    except deepl.DeepLException as e:
        logger.warning(f"Could not retrieve DeepL usage statistics: {e}")
        usage_str = _(
            "Current total API character usage: not available due to connection issues."
        )
    except Exception as e:
        logger.exception(f"Unexpected error retrieving usage stats: {e}")
        usage_str = _("Current total API character usage: not available.")

    return render(
        request,
        "statistics_frontend.html",
        {
            "usage": usage_str,
            "statistics": statistics,
            "granularity_title": granularity_title,
            "granularity": granularity,
            "dformat": dformat,
        },
    )


# ---------------------------------------------------------------------------
# Document translation (docx / pptx)
# ---------------------------------------------------------------------------

def _translate_text_fragment(translator, text, source_lang, target_lang,
                             formality, glossary, model_type, deadline=None):
    """Translate a single text fragment via the DeepL API.

    Returns the translated string, or the original text if the fragment
    is empty / whitespace-only.  Raises ``TimeoutError`` when *deadline*
    (a ``time.monotonic()`` value) has been reached.
    """
    if deadline and time.monotonic() > deadline:
        raise TimeoutError("Document translation wall-time limit exceeded.")
    if not text or not text.strip():
        return text
    result = translator.translate_text(
        text,
        source_lang=source_lang,
        target_lang=target_lang,
        formality=formality,
        glossary=glossary,
        model_type=model_type,
    )
    return result.text


def _translate_docx(filepath, translator, source_lang, target_lang,
                    formality, glossary, model_type, deadline=None,
                    count_only=False):
    """Translate all text in a .docx file **in-place**.

    Paragraphs from the document body, headers, footers, and tables are
    processed.  The full paragraph text is translated as a single unit
    (giving the translator proper sentence context) and placed into the
    first run.  Remaining runs are emptied so that the first run's
    character formatting is applied to the whole translated text.

    Word often splits visually identical text across many runs due to
    editing-history tracking (rsid attributes), spell-check state, etc.
    Translating per-run would (a) send tiny fragments to the API –
    degrading translation quality – and (b) lose inter-run whitespace
    because the API trims leading/trailing spaces from short fragments.

    When *count_only* is ``True`` the file is parsed but no API calls
    are made and the file is **not** modified.  This is used for the
    pre-flight character-count check and guarantees the exact same
    traversal logic as the real translation pass.

    Returns ``(char_count, api_calls)``.
    """
    from docx import Document

    doc = Document(filepath)
    char_count = 0
    api_calls = 0

    def process_paragraph(para):
        nonlocal char_count, api_calls
        runs = para.runs
        if not runs:
            return
        full_text = "".join(run.text for run in runs)
        if not full_text or not full_text.strip():
            return
        char_count += len(full_text)
        if count_only:
            return
        api_calls += 1
        translated = _translate_text_fragment(
            translator, full_text,
            source_lang, target_lang,
            formality, glossary, model_type,
            deadline=deadline,
        )
        # Place translated text in first run, clear the rest
        runs[0].text = translated
        for run in runs[1:]:
            run.text = ""

    # Body paragraphs
    for para in doc.paragraphs:
        process_paragraph(para)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    process_paragraph(para)

    # Headers and footers (deduplicate linked sections that share the
    # same XML element to avoid translating the same text twice).
    # We must keep strong references to the lxml element proxies;
    # otherwise the GC can collect them and a subsequent access to the
    # same XML node creates a *new* proxy with a different id().
    seen_elements = []
    seen_element_ids = set()
    for section in doc.sections:
        for header_footer in (section.header, section.footer):
            if header_footer is not None:
                for para in header_footer.paragraphs:
                    elem = para._element
                    elem_id = id(elem)
                    if elem_id in seen_element_ids:
                        continue
                    seen_element_ids.add(elem_id)
                    seen_elements.append(elem)  # prevent GC
                    process_paragraph(para)

    if not count_only:
        doc.save(filepath)
    return char_count, api_calls


def _translate_pptx(filepath, translator, source_lang, target_lang,
                    formality, glossary, model_type, deadline=None,
                    count_only=False):
    """Translate all text in a .pptx file **in-place**.

    Every text frame on every slide (including grouped shapes, tables and
    notes) is processed.  The full paragraph text is translated as a
    single unit and placed into the first run; remaining runs are emptied.
    See ``_translate_docx`` for the rationale.

    When *count_only* is ``True`` the file is parsed but no API calls
    are made and the file is **not** modified.

    Returns ``(char_count, api_calls)``.
    """
    from pptx import Presentation

    prs = Presentation(filepath)
    char_count = 0
    api_calls = 0

    def process_paragraph(para):
        nonlocal char_count, api_calls
        runs = para.runs
        if not runs:
            return
        full_text = "".join(run.text for run in runs)
        if not full_text or not full_text.strip():
            return
        char_count += len(full_text)
        if count_only:
            return
        api_calls += 1
        translated = _translate_text_fragment(
            translator, full_text,
            source_lang, target_lang,
            formality, glossary, model_type,
            deadline=deadline,
        )
        runs[0].text = translated
        for run in runs[1:]:
            run.text = ""

    def process_shape(shape):
        """Recursively process shapes, including groups and tables."""
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                process_paragraph(para)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        process_paragraph(para)
        if hasattr(shape, "shapes"):
            # Group shape – recurse into children
            for child in shape.shapes:
                process_shape(child)

    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape)

        # Slide notes
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                process_paragraph(para)

    if not count_only:
        prs.save(filepath)
    return char_count, api_calls


def _lang_suffix(target_lang):
    """Return a short language suffix for the output filename."""
    return target_lang[:2].lower()


# ---------------------------------------------------------------------------
# Document character counting (pre-flight check)
# ---------------------------------------------------------------------------
# Thin wrappers that reuse the translation functions in count-only mode.
# This guarantees the counting pass traverses exactly the same paragraphs
# as the real translation pass — no risk of diverging implementations.
# ---------------------------------------------------------------------------

def _count_chars_docx(filepath):
    """Return the total character count of a .docx file (without translating)."""
    char_count, _ = _translate_docx(
        filepath, None, None, None, None, None, None, count_only=True,
    )
    return char_count


def _count_chars_pptx(filepath):
    """Return the total character count of a .pptx file (without translating)."""
    char_count, _ = _translate_pptx(
        filepath, None, None, None, None, None, None, count_only=True,
    )
    return char_count


# ---------------------------------------------------------------------------
# Stale document job cleanup
# ---------------------------------------------------------------------------

_DOC_JOB_MAX_AGE_MINUTES = 10


def _cleanup_stale_document_jobs():
    """Delete translated files older than *_DOC_JOB_MAX_AGE_MINUTES* minutes.

    Called lazily when the translation form is loaded – no background
    process required.
    """
    cutoff = timezone.now() - timedelta(minutes=_DOC_JOB_MAX_AGE_MINUTES)
    stale = DocumentTranslationJob.objects.filter(
        created_at__lt=cutoff,
    ).exclude(result_path="")

    for job in stale:
        _delete_job_file(job)


def _delete_job_file(job):
    """Remove the translated file from disk (if it exists)."""
    if job.result_path:
        full = os.path.join(settings.MEDIA_ROOT, job.result_path)
        if os.path.isfile(full):
            try:
                os.remove(full)
            except OSError:
                pass
        # Also remove the parent directory if empty
        parent = os.path.dirname(full)
        try:
            os.rmdir(parent)
        except OSError:
            pass
        job.result_path = ""
        job.save(update_fields=["result_path"])


def _doc_upload_dir():
    """Return (and create) the base directory for translated documents."""
    d = os.path.join(settings.MEDIA_ROOT, "doc_translations")
    os.makedirs(d, exist_ok=True)
    return d


# ---------------------------------------------------------------------------
# Background translation worker
# ---------------------------------------------------------------------------

def _run_translation_job(job_id, language=None):
    """Execute the translation in a background thread.

    Reads the uploaded file from disk, translates it, and updates the
    ``DocumentTranslationJob`` row with the result metadata.
    """
    import django
    django.db.connections.close_all()
    if language:
        activate(language)

    try:
        job = DocumentTranslationJob.objects.get(pk=job_id)
    except DocumentTranslationJob.DoesNotExist:
        return

    job.status = DocumentTranslationJob.Status.PROCESSING
    job.save(update_fields=["status"])

    src_path = os.path.join(settings.MEDIA_ROOT, job.result_path)
    if not os.path.isfile(src_path):
        job.status = DocumentTranslationJob.Status.FAILED
        job.error_message = str(_("Upload file not found on disk."))
        job.save(update_fields=["status", "error_message"])
        return

    start = time.monotonic()
    timeout = getattr(settings, "DOCUMENT_TRANSLATION_TIMEOUT", 180)
    deadline = (start + timeout) if timeout > 0 else None

    # Pre-flight: count characters and check against MAX_TRANSLATION_LENGTH
    max_chars = getattr(settings, "MAX_TRANSLATION_LENGTH", 0)
    if max_chars > 0:
        try:
            if job.file_type == ".docx":
                doc_chars = _count_chars_docx(src_path)
            else:
                doc_chars = _count_chars_pptx(src_path)
        except Exception as e:
            logger.exception("Failed to count characters in document: %s", e)
            doc_chars = 0  # allow translation to proceed if counting fails

        if doc_chars > max_chars:
            exceeded_pct = round((doc_chars - max_chars) / max_chars * 100, 1)
            job.status = DocumentTranslationJob.Status.FAILED
            job.error_message = str(_(
                "The document contains %(doc_chars)s characters which exceeds "
                "the maximum allowed limit of %(max_chars)s characters "
                "(exceeded by %(exceeded_pct)s%%). Please use a shorter document."
            ) % {
                "doc_chars": f"{doc_chars:,}",
                "max_chars": f"{max_chars:,}",
                "exceeded_pct": exceeded_pct,
            })
            job.characters = doc_chars
            job.duration_seconds = round(time.monotonic() - start, 2)
            job.save(update_fields=[
                "status", "error_message", "characters", "duration_seconds",
            ])
            _delete_job_file(job)
            return

    try:
        translator = deepl.Translator(settings.DEEPL_AUTHKEY)

        # Parse direction
        direction_parts = job.direction.split("|")
        lang_pair = direction_parts[0].split("->")
        source_lang = lang_pair[0][:2]
        target_lang = lang_pair[1]
        formality = direction_parts[1] if len(direction_parts) > 1 and direction_parts[1] else None

        glossary_key = f"{source_lang}->{target_lang[:2]}"
        active_glossary = glossary_cache.get(glossary_key)
        model_type = DEEPL_MODEL_TYPE_TRANSLATION

        if job.file_type == ".docx":
            char_count, api_calls = _translate_docx(
                src_path, translator,
                source_lang, target_lang,
                formality, active_glossary, model_type,
                deadline=deadline,
            )
        else:
            char_count, api_calls = _translate_pptx(
                src_path, translator,
                source_lang, target_lang,
                formality, active_glossary, model_type,
                deadline=deadline,
            )

        elapsed = time.monotonic() - start

        # Log translation to statistics
        try:
            Translation.objects.create(
                characters=char_count,
                direction=job.direction,
                auto_detection=False,
                is_document_translation=True,
            )
        except Exception as e:
            logger.exception("Failed to log document translation: %s", e)

        job.characters = char_count
        job.api_calls = api_calls
        job.duration_seconds = round(elapsed, 2)
        job.status = DocumentTranslationJob.Status.COMPLETED
        job.completed_at = timezone.now()
        job.save(update_fields=[
            "characters", "api_calls", "duration_seconds",
            "status", "completed_at",
        ])

    except TimeoutError:
        logger.warning(
            "Document translation timed out after %s s (limit: %s s): %s",
            round(time.monotonic() - start, 2), timeout, job.id,
        )
        job.status = DocumentTranslationJob.Status.FAILED
        job.error_message = str(_(
            "The document translation exceeded the maximum allowed time "
            "of %(seconds)s seconds. Please try a smaller document."
        ) % {"seconds": timeout})
        job.duration_seconds = round(time.monotonic() - start, 2)
        job.save(update_fields=["status", "error_message", "duration_seconds"])
        _delete_job_file(job)
    except deepl.DeepLException as e:
        logger.error("DeepL API error during document translation: %s", e)
        job.status = DocumentTranslationJob.Status.FAILED
        job.error_message = str(_(
            "The translation service returned an error. "
            "Please try again later or check that the document "
            "is not corrupted."
        ))
        job.duration_seconds = round(time.monotonic() - start, 2)
        job.save(update_fields=["status", "error_message", "duration_seconds"])
        _delete_job_file(job)
    except ImportError as e:
        logger.error("Missing library for document translation: %s", e)
        job.status = DocumentTranslationJob.Status.FAILED
        job.error_message = str(_(
            "A required library for processing this file type is not installed. "
            "Please contact the administrator."
        ))
        job.duration_seconds = round(time.monotonic() - start, 2)
        job.save(update_fields=["status", "error_message", "duration_seconds"])
        _delete_job_file(job)
    except Exception as e:
        logger.exception("Unexpected error during document translation: %s", e)
        job.status = DocumentTranslationJob.Status.FAILED
        job.error_message = str(_(
            "An unexpected error occurred while translating the document. "
            "Please ensure the file is a valid .docx or .pptx document "
            "and try again."
        ))
        job.duration_seconds = round(time.monotonic() - start, 2)
        job.save(update_fields=["status", "error_message", "duration_seconds"])
        _delete_job_file(job)


# ---------------------------------------------------------------------------
# Document translation views
# ---------------------------------------------------------------------------

@require_POST
def deepl_document_translation(request):
    """Accept a document upload and start an async translation job.

    Returns a JSON response with the job ID for polling.
    """
    if not _is_doc_translation_enabled(request):
        return JsonResponse({"error": _("Document translation is not enabled.")}, status=404)

    form = DocumentTranslationForm(request.POST, request.FILES)
    if not form.is_valid():
        errors = {field: errs[0] for field, errs in form.errors.items()}
        return JsonResponse({"errors": errors}, status=400)

    # Ensure session exists (anonymous users get one too)
    if not request.session.session_key:
        request.session.create()

    # Reject if there is already a pending or processing job for this session
    active_job = DocumentTranslationJob.objects.filter(
        session_key=request.session.session_key,
        status__in=[
            DocumentTranslationJob.Status.PENDING,
            DocumentTranslationJob.Status.PROCESSING,
        ],
    ).first()
    if active_job:
        return JsonResponse(
            {"error": _("A translation is already in progress. Please wait for it to finish.")},
            status=409,
        )

    uploaded = form.cleaned_data["document"]
    direction = form.cleaned_data["directionChoice"]

    # Sanitise the filename to prevent path-traversal attacks.
    safe_name = os.path.basename(uploaded.name)
    if not safe_name:
        safe_name = "document" + os.path.splitext(uploaded.name)[1].lower()

    # Parse direction for the output filename
    direction_parts = direction.split("|")
    lang_pair = direction_parts[0].split("->")
    target_lang = lang_pair[1]
    ext = os.path.splitext(safe_name)[1].lower()
    base_name = os.path.splitext(safe_name)[0]
    out_filename = f"{base_name}_{_lang_suffix(target_lang)}{ext}"

    # Create the job record
    job = DocumentTranslationJob.objects.create(
        session_key=request.session.session_key,
        original_filename=safe_name,
        file_type=ext,
        file_size=uploaded.size,
        direction=direction,
        output_filename=out_filename,
    )

    # Save uploaded file to MEDIA_ROOT/doc_translations/<job_id>/
    job_dir = os.path.join(_doc_upload_dir(), str(job.id))
    os.makedirs(job_dir, exist_ok=True)
    file_path = os.path.join(job_dir, safe_name)
    with open(file_path, "wb") as f:
        for chunk in uploaded.chunks():
            f.write(chunk)

    job.result_path = os.path.join("doc_translations", str(job.id), safe_name)
    job.save(update_fields=["result_path"])

    # Start translation in background thread
    thread = threading.Thread(
        target=_run_translation_job,
        args=(job.id,),
        kwargs={"language": get_language()},
        daemon=True,
    )
    thread.start()

    return JsonResponse({"job_id": str(job.id)})


@require_GET
def deepl_document_job_status(request, job_id):
    """Return the current status of a translation job as JSON.

    Only the session that created the job may poll it.
    """
    if not _is_doc_translation_enabled(request):
        return JsonResponse({"error": _("Document translation is not enabled.")}, status=404)

    try:
        job = DocumentTranslationJob.objects.get(pk=job_id)
    except (DocumentTranslationJob.DoesNotExist, ValueError):
        return JsonResponse({"error": _("Job not found.")}, status=404)

    if job.session_key != request.session.session_key:
        return JsonResponse({"error": _("Access denied.")}, status=403)

    data = {
        "status": job.status,
        "original_filename": job.original_filename,
        "output_filename": job.output_filename,
    }
    if job.status == DocumentTranslationJob.Status.COMPLETED:
        data.update({
            "characters": job.characters,
            "api_calls": job.api_calls,
            "duration_seconds": job.duration_seconds,
            "file_type": job.file_type,
            "file_size": job.file_size,
        })
    elif job.status == DocumentTranslationJob.Status.FAILED:
        data["error_message"] = job.error_message

    return JsonResponse(data)


@require_GET
def deepl_document_download(request, job_id):
    """Serve the translated file and delete it from disk afterwards.

    Only the session that created the job may download it.
    """
    if not _is_doc_translation_enabled(request):
        return HttpResponse(_("Document translation is not enabled."), status=404)

    try:
        job = DocumentTranslationJob.objects.get(pk=job_id)
    except (DocumentTranslationJob.DoesNotExist, ValueError):
        return HttpResponse(_("Job not found."), status=404)

    if job.session_key != request.session.session_key:
        return HttpResponse(_("Access denied."), status=403)

    if job.status != DocumentTranslationJob.Status.COMPLETED:
        return HttpResponse(_("File not ready."), status=409)

    if not job.result_path:
        return HttpResponse(_("File no longer available on the server."), status=410)

    full_path = os.path.join(settings.MEDIA_ROOT, job.result_path)
    if not os.path.isfile(full_path):
        return HttpResponse(_("File no longer available on the server."), status=410)

    response = FileResponse(
        open(full_path, "rb"),
        as_attachment=True,
        filename=job.output_filename,
    )

    # Mark as downloaded; actual file cleanup happens via
    # _cleanup_stale_document_jobs() on the next form load.
    job.downloaded = True
    job.download_count = F("download_count") + 1
    job.save(update_fields=["downloaded", "download_count"])

    return response
