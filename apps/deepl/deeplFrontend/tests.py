# SPDX-License-Identifier: Apache-2.0
# SPDX-FileCopyrightText: 2024-2026 David Kleinhans, Jade University of Applied Sciences
"""
Test suite for DeepL Translation Frontend.

Tests cover:
- Model creation, validation, and properties
- Form validation and choices
- View functionality (GET requests, POST translations with mocked API)
- URL routing
- Context processor language selection and fallback chain
- Template tag functionality
- Settings configuration
"""

import os
import tempfile
from datetime import timedelta
from io import BytesIO
from unittest.mock import MagicMock, patch

from django.conf import settings
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import TestCase, Client, RequestFactory, override_settings
from django.urls import reverse
from django.utils import timezone

from .forms import TranslationForm, DocumentTranslationForm
from .models import Translation, Glossary, DocumentTranslationJob


# ============================================================================
# Model Tests
# ============================================================================

class TranslationModelTest(TestCase):
    """Test Translation model creation, defaults, and string representation."""

    def test_translation_creation(self):
        """Test creating a translation record with all fields."""
        t = Translation.objects.create(
            characters=150,
            direction="DE->EN-GB|",
            auto_detection=False,
        )
        self.assertEqual(t.characters, 150)
        self.assertEqual(t.direction, "DE->EN-GB|")
        self.assertFalse(t.auto_detection)
        self.assertIsNotNone(t.timestamp)

    def test_translation_str(self):
        """Test string representation includes direction and timestamp."""
        t = Translation.objects.create(
            characters=100,
            direction="EN-GB->DE|more",
            auto_detection=True,
        )
        s = str(t)
        self.assertIn("EN-GB->DE|more", s)

    def test_translation_ordering(self):
        """Test that translations are ordered by timestamp descending."""
        t1 = Translation.objects.create(
            characters=10, direction="DE->EN-GB|", auto_detection=False
        )
        t2 = Translation.objects.create(
            characters=20, direction="DE->EN-GB|", auto_detection=False
        )
        translations = list(Translation.objects.all())
        # Most recent first
        self.assertEqual(translations[0].pk, t2.pk)
        self.assertEqual(translations[1].pk, t1.pk)

    def test_translation_auto_detection_field(self):
        """Test auto_detection boolean field values."""
        t_auto = Translation.objects.create(
            characters=50, direction="DE->EN-GB|", auto_detection=True
        )
        t_manual = Translation.objects.create(
            characters=50, direction="DE->EN-GB|", auto_detection=False
        )
        self.assertTrue(t_auto.auto_detection)
        self.assertFalse(t_manual.auto_detection)


class GlossaryModelTest(TestCase):
    """Test Glossary model creation, properties, and string representation."""

    def setUp(self):
        """Create a test glossary."""
        self.glossary = Glossary.objects.create(
            glossary_id="test-glossary-id-123",
            name="Test Glossary",
            source_lang="DE",
            target_lang="EN-GB",
            original_filename="test_glossary.csv",
            comment="A test glossary",
            entry_count=42,
        )

    def test_glossary_creation(self):
        """Test glossary fields are stored correctly."""
        g = self.glossary
        self.assertEqual(g.glossary_id, "test-glossary-id-123")
        self.assertEqual(g.name, "Test Glossary")
        self.assertEqual(g.source_lang, "DE")
        self.assertEqual(g.target_lang, "EN-GB")
        self.assertEqual(g.original_filename, "test_glossary.csv")
        self.assertEqual(g.comment, "A test glossary")
        self.assertEqual(g.entry_count, 42)
        self.assertIsNotNone(g.upload_date)

    def test_glossary_str(self):
        """Test string representation format."""
        self.assertEqual(str(self.glossary), "Test Glossary (DE->EN-GB)")

    def test_glossary_language_pair(self):
        """Test language_pair property."""
        self.assertEqual(self.glossary.language_pair, "DE->EN-GB")

    def test_glossary_unique_id(self):
        """Test that glossary_id must be unique."""
        from django.db import IntegrityError

        with self.assertRaises(IntegrityError):
            Glossary.objects.create(
                glossary_id="test-glossary-id-123",  # Same as setUp
                name="Duplicate",
                source_lang="EN",
                target_lang="DE",
                original_filename="dup.csv",
            )

    def test_glossary_comment_blank(self):
        """Test that comment can be empty."""
        g = Glossary.objects.create(
            glossary_id="no-comment-id",
            name="No Comment Glossary",
            source_lang="EN",
            target_lang="DE",
            original_filename="nc.csv",
        )
        self.assertEqual(g.comment, "")

    def test_glossary_ordering(self):
        """Test that glossaries are ordered by upload_date descending."""
        g2 = Glossary.objects.create(
            glossary_id="second-glossary-id",
            name="Second Glossary",
            source_lang="EN",
            target_lang="DE",
            original_filename="second.csv",
        )
        glossaries = list(Glossary.objects.all())
        # Most recent first
        self.assertEqual(glossaries[0].pk, g2.pk)


# ============================================================================
# Form Tests
# ============================================================================

class TranslationFormTest(TestCase):
    """Test TranslationForm validation and field behavior."""

    def test_valid_form_auto(self):
        """Test form with auto detection direction."""
        form = TranslationForm(data={
            "sourceText": "Hello world",
            "directionChoice": "auto",
        })
        self.assertTrue(form.is_valid())

    def test_valid_form_de_en(self):
        """Test form with DE->EN direction."""
        form = TranslationForm(data={
            "sourceText": "Hallo Welt",
            "directionChoice": "DE->EN-GB|",
        })
        self.assertTrue(form.is_valid())

    def test_valid_form_en_de_formal(self):
        """Test form with EN->DE formal direction."""
        form = TranslationForm(data={
            "sourceText": "Hello world",
            "directionChoice": "EN-GB->DE|more",
        })
        self.assertTrue(form.is_valid())

    def test_valid_form_en_de_informal(self):
        """Test form with EN->DE informal direction."""
        form = TranslationForm(data={
            "sourceText": "Hello world",
            "directionChoice": "EN-GB->DE|less",
        })
        self.assertTrue(form.is_valid())

    def test_empty_source_text_invalid(self):
        """Test that empty source text is invalid."""
        form = TranslationForm(data={
            "sourceText": "",
            "directionChoice": "auto",
        })
        self.assertFalse(form.is_valid())
        self.assertIn("sourceText", form.errors)

    def test_missing_direction_invalid(self):
        """Test that missing direction choice is invalid."""
        form = TranslationForm(data={
            "sourceText": "Hello",
        })
        self.assertFalse(form.is_valid())
        self.assertIn("directionChoice", form.errors)

    def test_invalid_direction_invalid(self):
        """Test that an invalid direction choice is rejected."""
        form = TranslationForm(data={
            "sourceText": "Hello",
            "directionChoice": "INVALID",
        })
        self.assertFalse(form.is_valid())
        self.assertIn("directionChoice", form.errors)

    def test_direction_choices_count(self):
        """Test that exactly 4 direction choices are available."""
        form = TranslationForm()
        choices = form.fields["directionChoice"].choices
        self.assertEqual(len(choices), 4)

    def test_translated_text_disabled(self):
        """Test that translatedText field is disabled (output only)."""
        form = TranslationForm()
        self.assertTrue(form.fields["translatedText"].disabled)

    def test_translated_text_not_required(self):
        """Test that translatedText field is not required."""
        form = TranslationForm()
        self.assertFalse(form.fields["translatedText"].required)

    @override_settings(MAX_TRANSLATION_LENGTH=0)
    def test_no_length_limit_when_disabled(self):
        """Test that form accepts any length when MAX_TRANSLATION_LENGTH is 0."""
        long_text = "x" * 200000
        form = TranslationForm(data={
            "sourceText": long_text,
            "directionChoice": "auto",
        })
        self.assertTrue(form.is_valid())

    @override_settings(MAX_TRANSLATION_LENGTH=100)
    def test_form_valid_under_limit(self):
        """Test that text under the limit is accepted."""
        form = TranslationForm(data={
            "sourceText": "x" * 100,
            "directionChoice": "auto",
        })
        self.assertTrue(form.is_valid())

    @override_settings(MAX_TRANSLATION_LENGTH=100)
    def test_form_invalid_over_limit(self):
        """Test that text exceeding the limit is rejected."""
        form = TranslationForm(data={
            "sourceText": "x" * 101,
            "directionChoice": "auto",
        })
        self.assertFalse(form.is_valid())
        self.assertIn("sourceText", form.errors)

    @override_settings(MAX_TRANSLATION_LENGTH=100)
    def test_form_error_message_contains_limit(self):
        """Test that the error message mentions the character limit."""
        form = TranslationForm(data={
            "sourceText": "x" * 150,
            "directionChoice": "auto",
        })
        form.is_valid()
        error_message = form.errors["sourceText"][0]
        self.assertIn("100", error_message)
        self.assertIn("150", error_message)

    @override_settings(MAX_TRANSLATION_LENGTH=50)
    def test_form_exact_limit_is_accepted(self):
        """Test that text exactly at the limit is accepted."""
        form = TranslationForm(data={
            "sourceText": "x" * 50,
            "directionChoice": "auto",
        })
        self.assertTrue(form.is_valid())


# ============================================================================
# URL Tests
# ============================================================================

class URLRoutingTest(TestCase):
    """Test URL routing configuration."""

    def test_translation_form_url_resolves(self):
        """Test that translation-form URL resolves."""
        url = reverse("translation-form")
        self.assertIn("/deepl/translation/", url)

    def test_statistics_daily_url_resolves(self):
        """Test that daily statistics URL resolves."""
        url = reverse("usage-statistics", kwargs={"granularity": "d"})
        self.assertIn("stat/d", url)

    def test_statistics_monthly_url_resolves(self):
        """Test that monthly statistics URL resolves."""
        url = reverse("usage-statistics", kwargs={"granularity": "m"})
        self.assertIn("stat/m", url)

    def test_statistics_yearly_url_resolves(self):
        """Test that yearly statistics URL resolves."""
        url = reverse("usage-statistics", kwargs={"granularity": "y"})
        self.assertIn("stat/y", url)


# ============================================================================
# View Tests
# ============================================================================

class TranslationViewGetTest(TestCase):
    """Test GET requests to the translation view."""

    def setUp(self):
        self.client = Client()

    def test_translation_page_returns_200(self):
        """Test that translation page loads successfully."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.status_code, 200)

    def test_translation_page_contains_form(self):
        """Test that translation page contains the translation form."""
        response = self.client.get(reverse("translation-form"))
        self.assertContains(response, "sourceText")

    def test_translation_page_uses_correct_template(self):
        """Test that translation page uses the correct template."""
        response = self.client.get(reverse("translation-form"))
        self.assertTemplateUsed(response, "translation_frontend.html")

    def test_root_redirects_to_translation(self):
        """Test that root URL redirects to translation form."""
        response = self.client.get("/", follow=False)
        self.assertEqual(response.status_code, 301)

    @override_settings(MAX_TRANSLATION_LENGTH=500)
    def test_max_translation_length_in_context(self):
        """Test that MAX_TRANSLATION_LENGTH is passed to the template context."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.context["MAX_TRANSLATION_LENGTH"], 500)

    @override_settings(MAX_TRANSLATION_LENGTH=0)
    def test_max_translation_length_zero_in_context(self):
        """Test that MAX_TRANSLATION_LENGTH=0 (disabled) is in context."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.context["MAX_TRANSLATION_LENGTH"], 0)

    @override_settings(MAX_TRANSLATION_LENGTH=500)
    def test_char_limit_warning_in_page_when_enabled(self):
        """Test that the character limit warning div is present when limit is set."""
        response = self.client.get(reverse("translation-form"))
        self.assertContains(response, "char-limit-warning")

    @override_settings(MAX_TRANSLATION_LENGTH=0)
    def test_char_limit_warning_not_in_page_when_disabled(self):
        """Test that the character limit warning div is absent when limit is 0."""
        response = self.client.get(reverse("translation-form"))
        self.assertNotContains(response, "char-limit-warning")


class TranslationViewMaxLengthPostTest(TestCase):
    """Test POST requests when MAX_TRANSLATION_LENGTH is enforced."""

    def setUp(self):
        self.client = Client()

    @override_settings(MAX_TRANSLATION_LENGTH=100)
    def test_post_over_limit_rejected(self):
        """Test that a POST with text over the limit is rejected with form errors."""
        response = self.client.post(reverse("translation-form"), {
            "sourceText": "x" * 101,
            "directionChoice": "DE->EN-GB|",
        })
        self.assertEqual(response.status_code, 200)
        # The form should be re-rendered (not a redirect or API call)
        self.assertContains(response, "100")
        # No translation record should be created
        self.assertEqual(Translation.objects.count(), 0)

    @override_settings(MAX_TRANSLATION_LENGTH=100)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_under_limit_accepted(self, mock_translator_cls):
        """Test that a POST with text under the limit proceeds normally."""
        from .views import _GlossaryCache

        empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        empty_cache._cache = {}
        empty_cache._loaded_at = float("inf")

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with patch("deeplFrontend.views.glossary_cache", empty_cache):
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "x" * 100,
                "directionChoice": "DE->EN-GB|",
            })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Translated")
        self.assertEqual(Translation.objects.count(), 1)

    @override_settings(MAX_TRANSLATION_LENGTH=0)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_no_limit_when_disabled(self, mock_translator_cls):
        """Test that any length is accepted when MAX_TRANSLATION_LENGTH is 0."""
        from .views import _GlossaryCache

        empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        empty_cache._cache = {}
        empty_cache._loaded_at = float("inf")

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with patch("deeplFrontend.views.glossary_cache", empty_cache):
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "x" * 200000,
                "directionChoice": "DE->EN-GB|",
            })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Translated")


class TranslationViewPostTest(TestCase):
    """Test POST requests to the translation view with mocked DeepL API."""

    def setUp(self):
        self.client = Client()
        # Create a no-op glossary cache for tests that don't need glossaries
        from .views import _GlossaryCache

        self._empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        self._empty_cache._cache = {}
        self._empty_cache._loaded_at = float("inf")  # never expire

    def _patch_empty_glossary(self):
        """Return a context-manager that patches glossary_cache with an empty cache."""
        return patch("deeplFrontend.views.glossary_cache", self._empty_cache)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_de_en_translation(self, mock_translator_cls):
        """Test DE->EN translation request."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 1000
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "Hallo Welt",
                "directionChoice": "DE->EN-GB|",
            })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Hello World")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_en_de_formal_translation(self, mock_translator_cls):
        """Test EN->DE formal translation request."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Hallo Welt"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 1000
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "Hello World",
                "directionChoice": "EN-GB->DE|more",
            })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Hallo Welt")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_auto_detection_german(self, mock_translator_cls):
        """Test auto-detection when input is German."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        # First call: auto-detection (short text sent to determine language)
        mock_detect_result = MagicMock()
        mock_detect_result.detected_source_lang = "DE"
        # Second call: actual translation
        mock_translate_result = MagicMock()
        mock_translate_result.text = "Good morning"

        mock_translator.translate_text.side_effect = [
            mock_detect_result, mock_translate_result
        ]

        mock_usage = MagicMock()
        mock_usage.character.count = 1000
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "Guten Morgen",
                "directionChoice": "auto",
            })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Good morning")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_auto_detection_english(self, mock_translator_cls):
        """Test auto-detection when input is English."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_detect_result = MagicMock()
        mock_detect_result.detected_source_lang = "EN"
        mock_translate_result = MagicMock()
        mock_translate_result.text = "Guten Morgen"

        mock_translator.translate_text.side_effect = [
            mock_detect_result, mock_translate_result
        ]

        mock_usage = MagicMock()
        mock_usage.character.count = 1000
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "Good morning",
                "directionChoice": "auto",
            })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Guten Morgen")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_creates_translation_record(self, mock_translator_cls):
        """Test that a successful translation creates a database record."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Translated text"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        self.assertEqual(Translation.objects.count(), 0)
        with self._patch_empty_glossary():
            self.client.post(reverse("translation-form"), {
                "sourceText": "Test text",
                "directionChoice": "DE->EN-GB|",
            })
        self.assertEqual(Translation.objects.count(), 1)

        record = Translation.objects.first()
        self.assertEqual(record.direction, "DE->EN-GB|")
        self.assertEqual(record.characters, len("Test text"))
        self.assertFalse(record.auto_detection)
        self.assertFalse(record.is_document_translation)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_auto_detection_sets_flag(self, mock_translator_cls):
        """Test that auto-detection sets the auto_detection flag in DB."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_detect = MagicMock()
        mock_detect.detected_source_lang = "DE"
        mock_translate = MagicMock()
        mock_translate.text = "Result"

        mock_translator.translate_text.side_effect = [mock_detect, mock_translate]

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            self.client.post(reverse("translation-form"), {
                "sourceText": "Hallo",
                "directionChoice": "auto",
            })
        record = Translation.objects.first()
        self.assertTrue(record.auto_detection)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_api_error_returns_error_message(self, mock_translator_cls):
        """Test that DeepL API error returns user-friendly error message."""
        import deepl

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_translator.translate_text.side_effect = deepl.DeepLException("API Error")

        with self._patch_empty_glossary():
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "Test text",
                "directionChoice": "DE->EN-GB|",
            })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "ERROR")

    def test_post_invalid_form_returns_200(self):
        """Test that posting an invalid form still returns 200."""
        response = self.client.post(reverse("translation-form"), {
            "sourceText": "",
            "directionChoice": "auto",
        })
        self.assertEqual(response.status_code, 200)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_auto_detection_uses_latency_model(self, mock_translator_cls):
        """Test that auto-detection uses the latency-optimized model type."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_detect = MagicMock()
        mock_detect.detected_source_lang = "DE"
        mock_translate = MagicMock()
        mock_translate.text = "Translated"
        mock_translator.translate_text.side_effect = [mock_detect, mock_translate]

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            self.client.post(reverse("translation-form"), {
                "sourceText": "Hallo Welt",
                "directionChoice": "auto",
            })

        # First call is the detection call
        detection_call = mock_translator.translate_text.call_args_list[0]
        self.assertEqual(
            detection_call.kwargs.get("model_type"),
            "latency_optimized",
        )

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_translation_uses_configured_model(self, mock_translator_cls):
        """Test that the actual translation uses the configured model type."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            self.client.post(reverse("translation-form"), {
                "sourceText": "Hallo Welt",
                "directionChoice": "DE->EN-GB|",
            })

        from deeplFrontend.views import DEEPL_MODEL_TYPE_TRANSLATION
        call_kwargs = mock_translator.translate_text.call_args
        self.assertEqual(
            call_kwargs.kwargs.get("model_type"),
            DEEPL_MODEL_TYPE_TRANSLATION,
        )

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_auto_detection_translation_uses_configured_model(self, mock_translator_cls):
        """Test that after auto-detection, the actual translation uses the configured model."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_detect = MagicMock()
        mock_detect.detected_source_lang = "EN"
        mock_translate = MagicMock()
        mock_translate.text = "Guten Morgen"
        mock_translator.translate_text.side_effect = [mock_detect, mock_translate]

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            self.client.post(reverse("translation-form"), {
                "sourceText": "Good morning",
                "directionChoice": "auto",
            })

        # Second call is the actual translation
        from deeplFrontend.views import DEEPL_MODEL_TYPE_TRANSLATION
        translation_call = mock_translator.translate_text.call_args_list[1]
        self.assertEqual(
            translation_call.kwargs.get("model_type"),
            DEEPL_MODEL_TYPE_TRANSLATION,
        )

    @patch("deeplFrontend.views.deepl.Translator")
    def test_post_displays_usage_info(self, mock_translator_cls):
        """Test that successful translation displays API usage information."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 12345
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        with self._patch_empty_glossary():
            response = self.client.post(reverse("translation-form"), {
                "sourceText": "Test",
                "directionChoice": "DE->EN-GB|",
            })
        self.assertContains(response, "12345")


class StatisticsViewTest(TestCase):
    """Test statistics views with different granularities."""

    def setUp(self):
        self.client = Client()
        # Create some test translation records
        now = timezone.now()
        for i in range(3):
            Translation.objects.create(
                characters=100 + i * 50,
                direction="DE->EN-GB|",
                auto_detection=False,
            )
        Translation.objects.create(
            characters=200,
            direction="EN-GB->DE|more",
            auto_detection=True,
        )

    @patch("deeplFrontend.views.deepl.Translator")
    def test_daily_statistics_returns_200(self, mock_translator_cls):
        """Test daily statistics page loads."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 1000
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_monthly_statistics_returns_200(self, mock_translator_cls):
        """Test monthly statistics page loads."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 1000
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "m"})
        )
        self.assertEqual(response.status_code, 200)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_yearly_statistics_returns_200(self, mock_translator_cls):
        """Test yearly statistics page loads."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 1000
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "y"})
        )
        self.assertEqual(response.status_code, 200)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_statistics_uses_correct_template(self, mock_translator_cls):
        """Test statistics page uses correct template."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertTemplateUsed(response, "statistics_frontend.html")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_invalid_granularity_still_returns_200(self, mock_translator_cls):
        """Test that invalid granularity returns 200 with error message."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "x"})
        )
        self.assertEqual(response.status_code, 200)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_statistics_context_contains_expected_keys(self, mock_translator_cls):
        """Test that statistics context has all required keys."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 500
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertIn("statistics", response.context)
        self.assertIn("granularity_title", response.context)
        self.assertIn("usage", response.context)
        self.assertIn("dformat", response.context)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_statistics_api_error_handled(self, mock_translator_cls):
        """Test that API error during statistics is handled gracefully."""
        import deepl

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_translator.get_usage.side_effect = deepl.DeepLException("API Error")

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_statistics_doc_key_present_in_context(self, mock_translator_cls):
        """Test that each statistics period dict contains the 'doc' key."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        statistics = response.context["statistics"]
        self.assertTrue(len(statistics) > 0)
        for period in statistics:
            self.assertIn("doc", period, "Each period dict must contain a 'doc' key")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_statistics_doc_percentage_formatting(self, mock_translator_cls):
        """Test doc column shows 'N (X%)' format when document translations exist."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        # Add one document translation to the 4 already created in setUp (3+1=4 text).
        # Total for today: 5 translations, 1 is a document → 20 %.
        Translation.objects.create(
            characters=300,
            direction="DE->EN-GB|",
            auto_detection=False,
            is_document_translation=True,
        )

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        # The first period in the daily view covers today and must contain our records.
        today_period = response.context["statistics"][0]
        # 5 total, 1 doc → "1 (20%)"
        self.assertEqual(today_period["doc"], "1 (20%)")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_statistics_doc_zero_formatting(self, mock_translator_cls):
        """Test doc column shows '0' when no document translations exist."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        # All setUp translations have is_document_translation=False (default).
        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        today_period = response.context["statistics"][0]
        self.assertEqual(today_period["doc"], "0")

    @patch("deeplFrontend.views.deepl.Translator")
    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_statistics_doc_column_visible_when_enabled(self, mock_translator_cls):
        """Test that the doc column header is rendered when doc translation is enabled."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "whereof documents")

    @patch("deeplFrontend.views.deepl.Translator")
    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_statistics_doc_column_hidden_when_disabled(self, mock_translator_cls):
        """Test that the doc column is absent when doc translation is disabled."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        self.assertNotContains(response, "whereof documents")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_statistics_granularity_nav_present(self, mock_translator_cls):
        """Test that the granularity navigation pills are rendered on the statistics page."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, 'nav-pills')
        # Active pill must reflect current granularity
        self.assertContains(response, 'nav-link active')

    @patch("deeplFrontend.views.deepl.Translator")
    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_statistics_doc_column_renders_data_cells(self, mock_translator_cls):
        """Test that doc column <td> values are rendered when doc translation is enabled."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        # One doc translation exists (from setUp's 4 text translations + this one).
        Translation.objects.create(
            characters=150,
            direction="DE->EN-GB|",
            auto_detection=False,
            is_document_translation=True,
        )

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        # The formatted doc value for today must appear somewhere in the rendered table.
        today_period = response.context["statistics"][0]
        self.assertIn("(", today_period["doc"], "Doc value should contain a percentage")
        self.assertContains(response, today_period["doc"])

    @patch("deeplFrontend.views.deepl.Translator")
    @override_settings(DOCUMENT_TRANSLATION_ENABLED="@eligible\\.example$")
    def test_statistics_doc_column_visible_for_regex_eligible_user(self, mock_translator_cls):
        """Test doc column visible when setting is a regex and the user's session flag is set.

        This covers the tri-state 'specific user' eligibility path used in production
        deployments where DOCUMENT_TRANSLATION_ENABLED is an email-matching regex.
        """
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        # Simulate the session flag that DocumentTranslationsMiddleware sets for
        # users whose email matches the regex.
        session = self.client.session
        session["document_translations"] = True
        session.save()

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "whereof documents")

    @patch("deeplFrontend.views.deepl.Translator")
    @override_settings(DOCUMENT_TRANSLATION_ENABLED="@eligible\\.example$")
    def test_statistics_doc_column_hidden_for_regex_ineligible_user(self, mock_translator_cls):
        """Test doc column hidden when setting is a regex but the user has no session flag.

        Covers the case where another user on the same instance does NOT have access.
        """
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        # Ensure no session flag is present (ineligible user).
        session = self.client.session
        session.pop("document_translations", None)
        session.save()

        response = self.client.get(
            reverse("usage-statistics", kwargs={"granularity": "d"})
        )
        self.assertEqual(response.status_code, 200)
        self.assertNotContains(response, "whereof documents")


# ============================================================================
# Template Tag Tests
# ============================================================================

class LoadAboutContentTest(TestCase):
    """Test the load_about_content template tag."""

    def test_load_about_content_with_file(self):
        """Test loading about content from a markdown file."""
        from .templatetags.custom_helper_tags import load_about_content

        with tempfile.TemporaryDirectory() as tmpdir:
            # Create a test markdown file
            md_file = os.path.join(tmpdir, "about_en.md")
            with open(md_file, "w") as f:
                f.write("# Test About\n\nThis is a **test**.")

            with patch(
                "deeplFrontend.templatetags.custom_helper_tags.Path"
            ) as mock_path_cls:
                mock_path = MagicMock()
                mock_path_cls.return_value = mock_path

                # Make the path resolve to our temp file
                mock_md_file = MagicMock()
                mock_md_file.exists.return_value = True
                mock_md_file.__str__ = lambda self: md_file
                mock_md_file.__fspath__ = lambda self: md_file
                mock_path.__truediv__ = lambda self, name: mock_md_file

                # Mock open to read from the actual file
                result = load_about_content()
                # The function should return HTML content (may use default path)
                self.assertIsNotNone(result)


class ChangeLanguageTagTest(TestCase):
    """Test the change_lang template tag."""

    def test_change_lang_returns_url(self):
        """Test that change_lang returns a translated URL."""
        from .templatetags.custom_helper_tags import change_lang

        factory = RequestFactory()
        # Use a real resolvable URL so translate_url can match it
        url = reverse("translation-form")
        request = factory.get(url)
        context = {"request": request}

        result = change_lang(context, "de")
        self.assertIn("/de/", result)
        self.assertIn("/deepl/translation/", result)


# ============================================================================
# Settings Tests
# ============================================================================

class SettingsTest(TestCase):
    """Test that essential settings are configured correctly for tests."""

    def test_debug_is_false_during_tests(self):
        """Test that Django forces DEBUG=False during test runs for safety."""
        self.assertFalse(settings.DEBUG)

    def test_database_is_sqlite(self):
        """Test that test database uses SQLite."""
        self.assertIn("sqlite3", settings.DATABASES["default"]["ENGINE"])

    def test_deepl_authkey_configured(self):
        """Test that DEEPL_AUTHKEY is set (dummy for tests)."""
        self.assertTrue(hasattr(settings, "DEEPL_AUTHKEY"))
        self.assertTrue(len(settings.DEEPL_AUTHKEY) > 0)

    def test_translation_settings_have_defaults(self):
        """Test that translation settings have sensible defaults."""
        self.assertGreater(settings.TRANSLATION_DETECTION_LEN, 0)
        self.assertGreater(settings.STATISTICS_DAYS, 0)
        self.assertGreater(settings.STATISTICS_MONTHS, 0)
        self.assertGreater(settings.STATISTICS_YEARS, 0)

    def test_deepl_model_type_settings_exist(self):
        """Test that DeepL model type settings are configured."""
        self.assertTrue(hasattr(settings, "DEEPL_MODEL_TYPE_DETECTION"))
        self.assertTrue(hasattr(settings, "DEEPL_MODEL_TYPE_TRANSLATION"))
        self.assertIn(
            settings.DEEPL_MODEL_TYPE_DETECTION,
            ("", "quality_optimized", "prefer_quality_optimized", "latency_optimized"),
        )
        self.assertIn(
            settings.DEEPL_MODEL_TYPE_TRANSLATION,
            ("", "quality_optimized", "prefer_quality_optimized", "latency_optimized"),
        )

    def test_installed_apps_includes_app(self):
        """Test that deeplFrontend is in INSTALLED_APPS."""
        self.assertIn("deeplFrontend", settings.INSTALLED_APPS)

    def test_crispy_forms_configured(self):
        """Test that crispy forms is configured with Bootstrap 5."""
        self.assertEqual(settings.CRISPY_TEMPLATE_PACK, "bootstrap5")

    def test_app_title_is_dict(self):
        """Test that APP_TITLE is a dict with language keys."""
        self.assertIsInstance(settings.APP_TITLE, dict)
        self.assertIn("en", settings.APP_TITLE)
        self.assertIn("de", settings.APP_TITLE)

    def test_organization_name_is_dict(self):
        """Test that ORGANIZATION_NAME is a dict with language keys."""
        self.assertIsInstance(settings.ORGANIZATION_NAME, dict)
        self.assertIn("en", settings.ORGANIZATION_NAME)
        self.assertIn("de", settings.ORGANIZATION_NAME)

    def test_max_translation_length_setting_exists(self):
        """Test that MAX_TRANSLATION_LENGTH setting exists and defaults to 0."""
        self.assertTrue(hasattr(settings, "MAX_TRANSLATION_LENGTH"))
        self.assertIsInstance(settings.MAX_TRANSLATION_LENGTH, int)
        self.assertGreaterEqual(settings.MAX_TRANSLATION_LENGTH, 0)

    def test_document_translation_notice_is_dict(self):
        """Test that DOCUMENT_TRANSLATION_NOTICE is a dict with language keys."""
        self.assertIsInstance(settings.DOCUMENT_TRANSLATION_NOTICE, dict)
        self.assertIn("en", settings.DOCUMENT_TRANSLATION_NOTICE)
        self.assertIn("de", settings.DOCUMENT_TRANSLATION_NOTICE)
        # Both defaults must be non-empty (they serve as sensible fallbacks)
        self.assertTrue(settings.DOCUMENT_TRANSLATION_NOTICE["en"])
        self.assertTrue(settings.DOCUMENT_TRANSLATION_NOTICE["de"])

    def test_collect_i18n_env_false_sentinel_clears_default(self):
        """_collect_i18n_env: setting a language variant to 'False' removes the default."""
        from config.settings import _collect_i18n_env
        import unittest.mock as mock
        env = {"MY_SETTING_EN": "False", "MY_SETTING_DE": "Hallo"}
        with mock.patch.dict("os.environ", env, clear=False):
            result = _collect_i18n_env("MY_SETTING", {"en": "Hello", "de": "Hallo"})
        self.assertNotIn("en", result)
        self.assertIn("de", result)

    def test_collect_i18n_env_false_case_insensitive(self):
        """_collect_i18n_env: 'FALSE', 'false', 'False' all clear the default."""
        from config.settings import _collect_i18n_env
        import unittest.mock as mock
        for variant in ("FALSE", "false", "False", "fAlSe"):
            with self.subTest(variant=variant):
                env = {"MY_SETTING_EN": variant}
                with mock.patch.dict("os.environ", env, clear=False):
                    result = _collect_i18n_env("MY_SETTING", {"en": "Hello"})
                self.assertNotIn("en", result, f"Expected 'en' to be cleared for variant '{variant}'")

    def test_collect_i18n_env_all_false_produces_empty_dict(self):
        """_collect_i18n_env: setting all language variants to 'False' yields {}."""
        from config.settings import _collect_i18n_env
        import unittest.mock as mock
        env = {"MY_SETTING_EN": "False", "MY_SETTING_DE": "False"}
        with mock.patch.dict("os.environ", env, clear=False):
            result = _collect_i18n_env("MY_SETTING", {"en": "Hello", "de": "Hallo"})
        self.assertEqual(result, {})

    def test_collect_i18n_env_empty_string_still_ignored(self):
        """_collect_i18n_env: empty-string env vars are still silently ignored."""
        from config.settings import _collect_i18n_env
        import unittest.mock as mock
        env = {"MY_SETTING_EN": ""}
        with mock.patch.dict("os.environ", env, clear=False):
            result = _collect_i18n_env("MY_SETTING", {"en": "Hello"})
        # Empty string must NOT clear the default (use "False" for that)
        self.assertEqual(result.get("en"), "Hello")


# ============================================================================
# Context Processor Tests
# ============================================================================

class GetLocalisedTest(TestCase):
    """Test the _get_localised() fallback chain."""

    def test_returns_current_language(self):
        """Test that the current language value is returned."""
        from config.context_processors import _get_localised
        with self.settings(LANGUAGE_CODE="en-us"):
            with patch("config.context_processors.get_language", return_value="de"):
                result = _get_localised({"en": "English Title", "de": "Deutscher Titel"})
        self.assertEqual(result, "Deutscher Titel")

    def test_falls_back_to_language_code(self):
        """Test fallback to LANGUAGE_CODE when current language not in dict."""
        from config.context_processors import _get_localised
        with self.settings(LANGUAGE_CODE="en-us"):
            with patch("config.context_processors.get_language", return_value="fr"):
                result = _get_localised({"en": "English Title", "de": "Deutscher Titel"})
        self.assertEqual(result, "English Title")

    def test_falls_back_to_english(self):
        """Test fallback to 'en' when LANGUAGE_CODE not in dict either."""
        from config.context_processors import _get_localised
        with self.settings(LANGUAGE_CODE="fr"):
            with patch("config.context_processors.get_language", return_value="es"):
                result = _get_localised({"en": "English Title", "de": "Deutscher Titel"})
        self.assertEqual(result, "English Title")

    def test_falls_back_to_first_available(self):
        """Test fallback to first value when no standard key matches."""
        from config.context_processors import _get_localised
        with self.settings(LANGUAGE_CODE="fr"):
            with patch("config.context_processors.get_language", return_value="es"):
                result = _get_localised({"de": "Deutscher Titel"})
        self.assertEqual(result, "Deutscher Titel")

    def test_returns_fallback_string_for_empty_dict(self):
        """Test that empty dict returns the fallback string."""
        from config.context_processors import _get_localised
        result = _get_localised({}, fallback="Fallback")
        self.assertEqual(result, "Fallback")

    def test_handles_regional_language_code(self):
        """Test that regional codes like 'en-us' are normalised to 'en'."""
        from config.context_processors import _get_localised
        with self.settings(LANGUAGE_CODE="en-us"):
            with patch("config.context_processors.get_language", return_value="en-us"):
                result = _get_localised({"en": "English Title", "de": "Deutscher Titel"})
        self.assertEqual(result, "English Title")


class ContextProcessorTest(TestCase):
    """Test the app_settings context processor end-to-end."""

    def test_app_title_in_context(self):
        """Test that APP_TITLE appears in rendered pages."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.status_code, 200)
        self.assertIn("APP_TITLE", response.context)
        self.assertIsInstance(response.context["APP_TITLE"], str)
        self.assertTrue(len(response.context["APP_TITLE"]) > 0)

    def test_organization_name_in_context(self):
        """Test that ORGANIZATION_NAME appears in rendered pages."""
        response = self.client.get(reverse("translation-form"), follow=True)
        self.assertIn("ORGANIZATION_NAME", response.context)

    @override_settings(APP_TITLE={"en": "Test Title EN", "de": "Test Titel DE"})
    def test_app_title_language_selection(self):
        """Test that APP_TITLE uses the correct language."""
        # Request with German language
        response = self.client.get(reverse("translation-form"), HTTP_ACCEPT_LANGUAGE="de")
        # The title should be one of the configured values
        self.assertIn(response.context["APP_TITLE"], ["Test Title EN", "Test Titel DE"])

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_document_translation_notice_absent_when_enabled_for_all(self):
        """DOCUMENT_TRANSLATION_NOTICE must be empty when feature is on for everyone."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.context["DOCUMENT_TRANSLATION_NOTICE"], "")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_document_translation_notice_absent_when_disabled(self):
        """DOCUMENT_TRANSLATION_NOTICE must be empty when feature is off."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.context["DOCUMENT_TRANSLATION_NOTICE"], "")

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED="@",
        DOCUMENT_TRANSLATION_NOTICE={"en": "Restricted notice", "de": "Eingeschränkter Hinweis"},
    )
    def test_document_translation_notice_shown_in_regex_mode_with_access(self):
        """Notice is shown to users who have access in regex mode."""
        # Simulate middleware granting access
        session = self.client.session
        session["document_translations"] = True
        session.save()
        response = self.client.get(reverse("translation-form"))
        self.assertIn(response.context["DOCUMENT_TRANSLATION_NOTICE"], ["Restricted notice", "Eingeschränkter Hinweis"])

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED="@",
        DOCUMENT_TRANSLATION_NOTICE={"en": "Restricted notice", "de": "Eingeschränkter Hinweis"},
    )
    def test_document_translation_notice_absent_in_regex_mode_without_access(self):
        """Notice is NOT shown to users who do not have access in regex mode."""
        # No session flag — user not granted access
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.context["DOCUMENT_TRANSLATION_NOTICE"], "")

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED="@",
        DOCUMENT_TRANSLATION_NOTICE={"en": "", "de": ""},
    )
    def test_document_translation_notice_suppressed_when_empty_string(self):
        """Notice is suppressed when the operator sets it to an empty string."""
        session = self.client.session
        session["document_translations"] = True
        session.save()
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.context["DOCUMENT_TRANSLATION_NOTICE"], "")

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED=True,
        DOCUMENT_TRANSLATION_FAIR_USE_TEXT={"en": "Please confirm.", "de": "Bitte bestätigen."},
    )
    def test_fair_use_text_in_context_when_set(self):
        """DOCUMENT_TRANSLATION_FAIR_USE_TEXT context var is non-empty when configured."""
        response = self.client.get(reverse("translation-form"))
        self.assertIn("DOCUMENT_TRANSLATION_FAIR_USE_TEXT", response.context)
        self.assertTrue(response.context["DOCUMENT_TRANSLATION_FAIR_USE_TEXT"])

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED=True,
        DOCUMENT_TRANSLATION_FAIR_USE_TEXT={},
    )
    def test_fair_use_text_empty_in_context_when_not_set(self):
        """DOCUMENT_TRANSLATION_FAIR_USE_TEXT context var is empty string when dict is {}."""
        response = self.client.get(reverse("translation-form"))
        self.assertIn("DOCUMENT_TRANSLATION_FAIR_USE_TEXT", response.context)
        self.assertEqual(response.context["DOCUMENT_TRANSLATION_FAIR_USE_TEXT"], "")


# ============================================================================
# Glossary View Integration Tests
# ============================================================================

class GlossaryCacheTest(TestCase):
    """Test the _GlossaryCache lazy-loading and auto-refresh mechanism."""

    def setUp(self):
        self.glossary = Glossary.objects.create(
            glossary_id="cache-test-id",
            name="Cache Test Glossary",
            source_lang="DE",
            target_lang="EN-GB",
            original_filename="cache.csv",
            entry_count=5,
        )

    @patch("deeplFrontend.views.deepl.Translator")
    def test_lazy_load_on_first_access(self, mock_translator_cls):
        """Cache is populated lazily on the first .get() call, not at creation."""
        from .views import _GlossaryCache

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_deepl_glossary = MagicMock()
        mock_translator.get_glossary.return_value = mock_deepl_glossary

        cache = _GlossaryCache()
        # Cache should NOT have called the API yet
        mock_translator_cls.assert_not_called()

        # First .get() triggers loading
        result = cache.get("DE->EN")
        self.assertEqual(result, mock_deepl_glossary)
        mock_translator.get_glossary.assert_called_once_with("cache-test-id")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_normalises_keys_to_two_letter_codes(self, mock_translator_cls):
        """Cache keys are normalised to 2-letter codes (e.g. DE->EN, not DE->EN-GB)."""
        from .views import _GlossaryCache

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_translator.get_glossary.return_value = MagicMock()

        cache = _GlossaryCache()
        self.assertIsNotNone(cache.get("DE->EN"))
        self.assertIsNone(cache.get("DE->EN-GB"))

    @patch("deeplFrontend.views.deepl.Translator")
    def test_cache_ttl_refresh(self, mock_translator_cls):
        """Cache refreshes after TTL expires."""
        from .views import _GlossaryCache

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        first_glossary = MagicMock(name="first")
        second_glossary = MagicMock(name="second")
        mock_translator.get_glossary.side_effect = [first_glossary, second_glossary]

        cache = _GlossaryCache()
        self.assertEqual(cache.get("DE->EN"), first_glossary)

        # Simulate TTL expiry (float('-inf') guarantees reload on any system)
        cache._loaded_at = float('-inf')
        self.assertEqual(cache.get("DE->EN"), second_glossary)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_reload_forces_refresh(self, mock_translator_cls):
        """Manual reload() forces an immediate refresh."""
        from .views import _GlossaryCache

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        first_glossary = MagicMock(name="first")
        second_glossary = MagicMock(name="second")
        mock_translator.get_glossary.side_effect = [first_glossary, second_glossary]

        cache = _GlossaryCache()
        cache.get("DE->EN")  # initial load
        cache.reload()       # forced refresh
        self.assertEqual(cache.get("DE->EN"), second_glossary)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_loaded_pairs_property(self, mock_translator_cls):
        """loaded_pairs returns the set of cached language-pair keys."""
        from .views import _GlossaryCache

        Glossary.objects.create(
            glossary_id="en-de-id", name="EN-DE test",
            source_lang="EN", target_lang="DE",
            original_filename="b.csv", entry_count=3,
        )

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_translator.get_glossary.return_value = MagicMock()

        cache = _GlossaryCache()
        self.assertEqual(cache.loaded_pairs, {"DE->EN", "EN->DE"})

    @patch("deeplFrontend.views.deepl.Translator")
    def test_api_failure_returns_empty_cache(self, mock_translator_cls):
        """If the DeepL API is unreachable the cache stays empty (no crash)."""
        from .views import _GlossaryCache

        mock_translator_cls.side_effect = Exception("API unreachable")

        cache = _GlossaryCache()
        self.assertIsNone(cache.get("DE->EN"))
        self.assertEqual(cache.loaded_pairs, set())

    @patch("deeplFrontend.views.deepl.Translator")
    def test_duplicate_pair_uses_newest_glossary(self, mock_translator_cls):
        """When multiple glossaries exist for the same language pair,
        the most recently uploaded one is used (others are skipped)."""
        from .views import _GlossaryCache

        # Create an older duplicate glossary for the same pair (DE->EN).
        # setUp already created one with glossary_id="cache-test-id".
        older = Glossary.objects.create(
            glossary_id="older-duplicate-id",
            name="Older Duplicate Glossary",
            source_lang="DE",
            target_lang="EN",
            original_filename="old.csv",
            entry_count=3,
        )
        # auto_now_add ignores explicit values, so backdate via update()
        Glossary.objects.filter(pk=older.pk).update(
            upload_date=self.glossary.upload_date - timedelta(days=1),
        )

        newest_deepl = MagicMock(name="newest")
        older_deepl = MagicMock(name="older")

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        # Glossary.objects.all() returns newest first (-upload_date).
        # The cache should keep only the first (newest) one it encounters.
        mock_translator.get_glossary.side_effect = (
            lambda gid: newest_deepl if gid == "cache-test-id" else older_deepl
        )

        cache = _GlossaryCache()
        result = cache.get("DE->EN")

        self.assertEqual(result, newest_deepl, "Should use the newest glossary")
        self.assertEqual(cache.loaded_pairs, {"DE->EN"})
        # Only one entry for DE->EN, not two
        self.assertEqual(len(cache._cache), 1)


class GlossaryIntegrationTest(TestCase):
    """Test glossary integration in the translation flow."""

    def setUp(self):
        self.client = Client()
        # Create a test glossary in the database
        self.glossary = Glossary.objects.create(
            glossary_id="integration-test-id",
            name="Integration Test Glossary",
            source_lang="DE",
            target_lang="EN-GB",
            original_filename="integration.csv",
            entry_count=10,
        )

    def _make_mock_glossary_cache(self, pairs):
        """Create a mock _GlossaryCache whose .get() uses the given dict."""
        from .views import _GlossaryCache

        cache = _GlossaryCache.__new__(_GlossaryCache)
        cache._cache = pairs
        cache._loaded_at = float("inf")  # never expire
        return cache

    @patch("deeplFrontend.views.deepl.Translator")
    def test_glossary_passed_to_translator(self, mock_translator_cls):
        """Test that glossary is passed to translator when available."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        # Cache key uses normalised 2-letter codes (DE->EN, not DE->EN-GB)
        test_glossary_obj = MagicMock()
        mock_cache = self._make_mock_glossary_cache({"DE->EN": test_glossary_obj})
        with patch("deeplFrontend.views.glossary_cache", mock_cache):
            self.client.post(reverse("translation-form"), {
                "sourceText": "Testtext",
                "directionChoice": "DE->EN-GB|",
            })

        # Verify translate_text was called with the glossary
        call_kwargs = mock_translator.translate_text.call_args
        self.assertEqual(call_kwargs.kwargs.get("glossary"), test_glossary_obj)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_glossary_passed_for_en_de_direction(self, mock_translator_cls):
        """Test that glossary is passed for EN->DE direction."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Übersetzt"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        test_glossary_obj = MagicMock()
        mock_cache = self._make_mock_glossary_cache({"EN->DE": test_glossary_obj})
        with patch("deeplFrontend.views.glossary_cache", mock_cache):
            self.client.post(reverse("translation-form"), {
                "sourceText": "Hello world",
                "directionChoice": "EN-GB->DE|more",
            })

        call_kwargs = mock_translator.translate_text.call_args
        self.assertEqual(call_kwargs.kwargs.get("glossary"), test_glossary_obj)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_no_glossary_when_not_available(self, mock_translator_cls):
        """Test that None is passed when no glossary matches."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 0
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        mock_cache = self._make_mock_glossary_cache({})
        with patch("deeplFrontend.views.glossary_cache", mock_cache):
            self.client.post(reverse("translation-form"), {
                "sourceText": "Hello",
                "directionChoice": "EN-GB->DE|more",
            })

        call_kwargs = mock_translator.translate_text.call_args
        self.assertIsNone(call_kwargs.kwargs.get("glossary"))

    @patch("deeplFrontend.views.deepl.Translator")
    def test_cache_normalises_keys(self, mock_translator_cls):
        """Test that _GlossaryCache normalises keys to 2-letter codes.

        This is the root cause of the glossary-not-used bug: glossaries
        uploaded with target_lang='EN-GB' were stored under 'DE->EN-GB' in the
        cache, but the translation view looked them up as 'DE->EN'
        (normalised to 2-letter codes).  After the fix both sides normalise.
        """
        from .views import _GlossaryCache

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        # Simulate a glossary fetched from DeepL API
        mock_deepl_glossary = MagicMock()
        mock_translator.get_glossary.return_value = mock_deepl_glossary

        # DB has glossary with target_lang "EN-GB" (regional variant)
        # _GlossaryCache should normalise to "DE->EN"
        cache = _GlossaryCache()
        self.assertIsNotNone(cache.get("DE->EN"))
        self.assertIsNone(cache.get("DE->EN-GB"))
        self.assertEqual(cache.get("DE->EN"), mock_deepl_glossary)


# ============================================================================
# Document Translation Form Tests
# ============================================================================

def _make_docx_bytes():
    """Create a minimal valid .docx file in memory and return its bytes."""
    from docx import Document
    buf = BytesIO()
    doc = Document()
    doc.add_paragraph("Hallo Welt")
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes():
    """Create a minimal valid .pptx file in memory and return its bytes."""
    from pptx import Presentation
    from pptx.util import Inches
    buf = BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Titel"
    slide.placeholders[1].text = "Hallo Welt"
    prs.save(buf)
    return buf.getvalue()


class DocumentTranslationFormTest(TestCase):
    """Test DocumentTranslationForm validation and field behaviour."""

    def test_no_auto_detection_choice(self):
        """Document form must NOT offer 'auto' as a direction choice."""
        form = DocumentTranslationForm()
        choice_values = [v for v, _ in form.fields["directionChoice"].choices]
        self.assertNotIn("auto", choice_values)

    def test_direction_choices_count(self):
        """3 direction choices plus 1 blank placeholder."""
        form = DocumentTranslationForm()
        self.assertEqual(len(form.fields["directionChoice"].choices), 4)
        # First choice is the blank placeholder
        self.assertEqual(form.fields["directionChoice"].choices[0][0], "")

    def test_valid_docx_upload(self):
        """Valid .docx file and direction should pass validation."""
        docx = SimpleUploadedFile("test.docx", _make_docx_bytes(),
                                  content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        form = DocumentTranslationForm(
            data={
                "directionChoice": "DE->EN-GB|",
                "fair_use_confirmed": "on",
            },
            files={"document": docx},
        )
        self.assertTrue(form.is_valid())

    def test_valid_pptx_upload(self):
        """Valid .pptx file and direction should pass validation."""
        pptx = SimpleUploadedFile("slides.pptx", _make_pptx_bytes(),
                                  content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        form = DocumentTranslationForm(
            data={
                "directionChoice": "EN-GB->DE|more",
                "fair_use_confirmed": "on",
            },
            files={"document": pptx},
        )
        self.assertTrue(form.is_valid())

    def test_invalid_extension_rejected(self):
        """Uploading a .txt file should be rejected."""
        txt = SimpleUploadedFile("readme.txt", b"Hello world", content_type="text/plain")
        form = DocumentTranslationForm(
            data={
                "directionChoice": "DE->EN-GB|",
                "fair_use_confirmed": "on",
            },
            files={"document": txt},
        )
        self.assertFalse(form.is_valid())
        self.assertIn("document", form.errors)

    def test_pdf_rejected(self):
        """Uploading a .pdf file should be rejected."""
        pdf = SimpleUploadedFile("report.pdf", b"%PDF-1.4", content_type="application/pdf")
        form = DocumentTranslationForm(
            data={
                "directionChoice": "DE->EN-GB|",
                "fair_use_confirmed": "on",
            },
            files={"document": pdf},
        )
        self.assertFalse(form.is_valid())
        self.assertIn("document", form.errors)

    def test_old_doc_format_rejected(self):
        """Uploading a .doc (legacy) file should be rejected."""
        doc = SimpleUploadedFile("old.doc", b"\xd0\xcf\x11", content_type="application/msword")
        form = DocumentTranslationForm(
            data={
                "directionChoice": "DE->EN-GB|",
                "fair_use_confirmed": "on",
            },
            files={"document": doc},
        )
        self.assertFalse(form.is_valid())
        self.assertIn("document", form.errors)

    def test_missing_file_rejected(self):
        """Submitting without a file should be rejected."""
        form = DocumentTranslationForm(
            data={
                "directionChoice": "DE->EN-GB|",
                "fair_use_confirmed": "on",
            },
            files={},
        )
        self.assertFalse(form.is_valid())
        self.assertIn("document", form.errors)

    def test_missing_direction_rejected(self):
        """Submitting without a direction choice should be rejected."""
        docx = SimpleUploadedFile("test.docx", _make_docx_bytes())
        form = DocumentTranslationForm(
            data={"fair_use_confirmed": "on"},
            files={"document": docx},
        )
        self.assertFalse(form.is_valid())
        self.assertIn("directionChoice", form.errors)

    def test_invalid_direction_rejected(self):
        """Submitting with an invalid direction choice should be rejected."""
        docx = SimpleUploadedFile("test.docx", _make_docx_bytes())
        form = DocumentTranslationForm(
            data={
                "directionChoice": "auto",
                "fair_use_confirmed": "on",
            },
            files={"document": docx},
        )
        self.assertFalse(form.is_valid())
        self.assertIn("directionChoice", form.errors)

    def test_file_too_large_rejected(self):
        """File exceeding 50 MB should be rejected."""
        # Create a file object that reports a large size without allocating memory
        big = SimpleUploadedFile("huge.docx", b"x", content_type="application/octet-stream")
        big.size = 51 * 1024 * 1024  # 51 MB
        form = DocumentTranslationForm(
            data={
                "directionChoice": "DE->EN-GB|",
                "fair_use_confirmed": "on",
            },
            files={"document": big},
        )
        self.assertFalse(form.is_valid())
        self.assertIn("document", form.errors)

    def test_direction_required(self):
        """directionChoice field is required."""
        form = DocumentTranslationForm()
        self.assertTrue(form.fields["directionChoice"].required)

    def test_document_required(self):
        """document field is required."""
        form = DocumentTranslationForm()
        self.assertTrue(form.fields["document"].required)

    def test_document_widget_accept_attr(self):
        """document field widget must carry an accept attribute covering .docx and .pptx."""
        form = DocumentTranslationForm()
        accept = form.fields["document"].widget.attrs.get("accept", "")
        self.assertIn(".docx", accept)
        self.assertIn(".pptx", accept)

    def test_fair_use_field_required_by_default(self):
        """fair_use_confirmed field is required when text is configured."""
        form = DocumentTranslationForm()
        self.assertTrue(form.fields["fair_use_confirmed"].required)

    @override_settings(DOCUMENT_TRANSLATION_FAIR_USE_TEXT={})
    def test_fair_use_form_valid_when_text_empty_and_on_submitted(self):
        """When DOCUMENT_TRANSLATION_FAIR_USE_TEXT is empty the hidden input
        sends 'on', which satisfies BooleanField(required=True)."""
        docx = SimpleUploadedFile("test.docx", _make_docx_bytes())
        form = DocumentTranslationForm(
            data={
                "directionChoice": "DE->EN-GB|",
                "fair_use_confirmed": "on",
            },
            files={"document": docx},
        )
        self.assertTrue(form.is_valid())


# ============================================================================
# Document Translation URL Tests
# ============================================================================

class DocumentTranslationURLTest(TestCase):
    """Test URL routing for document translation."""

    def test_document_translation_url_resolves(self):
        """document-translation URL should resolve."""
        url = reverse("document-translation")
        self.assertIn("translation/document", url)


# ============================================================================
# Document Translation View Tests
# ============================================================================

class DocumentTranslationViewGetTest(TestCase):
    """Test GET requests and feature-flag behaviour for document translation."""

    def setUp(self):
        self.client = Client()

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_tabs_shown_when_enabled(self):
        """Translation page should show tabs when feature is enabled."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "translationTabs")
        self.assertContains(response, "document-pane")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_notice_banner_absent_when_enabled_for_all(self):
        """Info notice must not appear in the HTML when the feature is on for everyone."""
        response = self.client.get(reverse("translation-form"))
        # Context value must be empty; no notice markup in page
        self.assertEqual(response.context["DOCUMENT_TRANSLATION_NOTICE"], "")
        self.assertNotContains(response, "Doc translation notice")

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED="@",
        DOCUMENT_TRANSLATION_NOTICE={"en": "Doc translation notice"},
    )
    def test_notice_banner_rendered_for_regex_access_user(self):
        """Info notice must appear in the HTML for a user who has regex-mode access."""
        session = self.client.session
        session["document_translations"] = True
        session.save()
        response = self.client.get(reverse("translation-form"))
        self.assertContains(response, "Doc translation notice")
        self.assertContains(response, "alert-info")

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED="@",
        DOCUMENT_TRANSLATION_NOTICE={"en": "Doc translation notice"},
    )
    def test_notice_banner_absent_for_regex_no_access_user(self):
        """Info banner must NOT appear for a user who does not have regex-mode access."""
        # No session flag — feature tab is hidden, no banner either
        response = self.client.get(reverse("translation-form"))
        self.assertNotContains(response, "Doc translation notice")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_tabs_hidden_when_disabled(self):
        """Translation page should NOT show tabs when feature is disabled."""
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.status_code, 200)
        self.assertNotContains(response, "translationTabs")
        self.assertNotContains(response, "document-pane")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_doc_form_in_context(self):
        """doc_form should be present in context when enabled."""
        response = self.client.get(reverse("translation-form"))
        self.assertIn("doc_form", response.context)
        self.assertIsInstance(response.context["doc_form"], DocumentTranslationForm)

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED=True,
        DOCUMENT_TRANSLATION_FAIR_USE_TEXT={"en": "I confirm.", "de": "Ich bestätige."},
    )
    def test_fair_use_checkbox_rendered_when_text_configured(self):
        """Visible checkbox is rendered when fair-use text is non-empty."""
        response = self.client.get(reverse("translation-form"))
        self.assertContains(response, 'id="id_fair_use_confirmed"')
        self.assertNotContains(response, 'type="hidden" name="fair_use_confirmed"')

    @override_settings(
        DOCUMENT_TRANSLATION_ENABLED=True,
        DOCUMENT_TRANSLATION_FAIR_USE_TEXT={},
    )
    def test_fair_use_hidden_input_rendered_when_text_empty(self):
        """Hidden input is rendered instead of checkbox when fair-use text is empty."""
        response = self.client.get(reverse("translation-form"))
        self.assertContains(response, 'type="hidden" name="fair_use_confirmed" value="on"')
        self.assertNotContains(response, 'id="id_fair_use_confirmed"')

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_get_document_url_returns_405_when_disabled(self):
        """GET to document URL should return 405 (POST only)."""
        response = self.client.get(reverse("document-translation"))
        self.assertEqual(response.status_code, 405)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_post_document_url_returns_404_when_disabled(self):
        """POST to document URL should return 404 when feature is disabled."""
        docx = SimpleUploadedFile("test.docx", _make_docx_bytes())
        response = self.client.post(reverse("document-translation"), {
            "directionChoice": "DE->EN-GB|",
            "document": docx,
            "fair_use_confirmed": "on",
        })
        self.assertEqual(response.status_code, 404)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_get_document_url_returns_405(self):
        """GET to document translation endpoint should return 405 (POST only)."""
        response = self.client.get(reverse("document-translation"))
        self.assertEqual(response.status_code, 405)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_active_job_shown_on_page_load(self):
        """If a PENDING job exists for the session, active_doc_job_json is in context."""
        # Create a session first
        self.client.get(reverse("translation-form"))
        session = self.client.session
        session.save()
        job = DocumentTranslationJob.objects.create(
            session_key=session.session_key,
            original_filename="test.docx",
            file_type=".docx",
            file_size=1024,
            direction="DE->EN-GB|",
            output_filename="test_en.docx",
            status=DocumentTranslationJob.Status.PENDING,
        )
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.status_code, 200)
        self.assertIn("active_doc_job_json", response.context)
        import json
        data = json.loads(response.context["active_doc_job_json"])
        self.assertEqual(data["id"], str(job.id))
        self.assertEqual(data["status"], "pending")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_completed_not_downloaded_shown_on_reload(self):
        """A completed job that has NOT been downloaded should appear on reload."""
        self.client.get(reverse("translation-form"))
        session = self.client.session
        session.save()
        job = DocumentTranslationJob.objects.create(
            session_key=session.session_key,
            original_filename="test.docx",
            file_type=".docx",
            file_size=1024,
            direction="DE->EN-GB|",
            output_filename="test_en.docx",
            status=DocumentTranslationJob.Status.COMPLETED,
            result_path="doc_translations/fake/test.docx",
            downloaded=False,
            characters=500,
            duration_seconds=2.5,
        )
        response = self.client.get(reverse("translation-form"))
        self.assertEqual(response.status_code, 200)
        self.assertIn("active_doc_job_json", response.context)
        import json
        data = json.loads(response.context["active_doc_job_json"])
        self.assertEqual(data["id"], str(job.id))
        self.assertEqual(data["status"], "completed")
        self.assertEqual(data["characters"], 500)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_no_active_job_when_completed_and_cleaned(self):
        """A completed job with empty result_path should not appear as active."""
        self.client.get(reverse("translation-form"))
        session = self.client.session
        session.save()
        DocumentTranslationJob.objects.create(
            session_key=session.session_key,
            original_filename="test.docx",
            file_type=".docx",
            file_size=1024,
            direction="DE->EN-GB|",
            output_filename="test_en.docx",
            status=DocumentTranslationJob.Status.COMPLETED,
            result_path="",  # cleaned up
        )
        response = self.client.get(reverse("translation-form"))
        self.assertNotIn("active_doc_job_json", response.context)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_no_active_job_after_download(self):
        """A completed+downloaded job should not appear on page reload."""
        self.client.get(reverse("translation-form"))
        session = self.client.session
        session.save()
        DocumentTranslationJob.objects.create(
            session_key=session.session_key,
            original_filename="test.docx",
            file_type=".docx",
            file_size=1024,
            direction="DE->EN-GB|",
            output_filename="test_en.docx",
            status=DocumentTranslationJob.Status.COMPLETED,
            result_path="doc_translations/fake/test.docx",
            downloaded=True,
            download_count=1,
        )
        response = self.client.get(reverse("translation-form"))
        self.assertNotIn("active_doc_job_json", response.context)


class DocumentTranslationViewPostTest(TestCase):
    """Test POST requests to the document translation view with mocked DeepL API."""

    def setUp(self):
        self.client = Client()
        from .views import _GlossaryCache
        self._empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        self._empty_cache._cache = {}
        self._empty_cache._loaded_at = float("inf")

    def _patch_empty_glossary(self):
        return patch("deeplFrontend.views.glossary_cache", self._empty_cache)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_docx_translation_returns_job_id(self, mock_translator_cls):
        """Uploading a valid .docx should return JSON with a job_id."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        docx_bytes = _make_docx_bytes()
        docx = SimpleUploadedFile("bericht.docx", docx_bytes,
                                  content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread") as mock_thread:
            mock_thread.return_value = MagicMock()
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        data = response.json()
        self.assertIn("job_id", data)
        job = DocumentTranslationJob.objects.get(pk=data["job_id"])
        self.assertEqual(job.original_filename, "bericht.docx")
        self.assertEqual(job.output_filename, "bericht_en.docx")
        self.assertEqual(job.file_type, ".docx")
        mock_thread.return_value.start.assert_called_once()

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_pptx_translation_returns_job_id(self, mock_translator_cls):
        """Uploading a valid .pptx should return JSON with a job_id."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        pptx_bytes = _make_pptx_bytes()
        pptx = SimpleUploadedFile("folien.pptx", pptx_bytes,
                                  content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread") as mock_thread:
            mock_thread.return_value = MagicMock()
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": pptx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        data = response.json()
        self.assertIn("job_id", data)
        job = DocumentTranslationJob.objects.get(pk=data["job_id"])
        self.assertEqual(job.original_filename, "folien.pptx")
        self.assertEqual(job.output_filename, "folien_en.pptx")
        self.assertEqual(job.file_type, ".pptx")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_output_filename_has_target_lang_suffix(self, mock_translator_cls):
        """Output file should be named <original>_<lang>.<ext>."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Hallo"
        mock_translator.translate_text.return_value = mock_result

        docx = SimpleUploadedFile("report.docx", _make_docx_bytes())

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread") as mock_thread:
            mock_thread.return_value = MagicMock()
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "EN-GB->DE|more",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        data = response.json()
        job = DocumentTranslationJob.objects.get(pk=data["job_id"])
        self.assertEqual(job.output_filename, "report_de.docx")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_translation_record_created(self, mock_translator_cls):
        """Background translation should create a Translation record."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        self.assertEqual(Translation.objects.count(), 0)
        docx = SimpleUploadedFile("test.docx", _make_docx_bytes())

        # Use a synchronous thread mock that runs the target immediately
        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(Translation.objects.count(), 1)
        record = Translation.objects.first()
        self.assertEqual(record.direction, "DE->EN-GB|")
        self.assertFalse(record.auto_detection)
        self.assertGreater(record.characters, 0)
        self.assertTrue(record.is_document_translation)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_api_error_sets_job_failed(self, mock_translator_cls):
        """DeepL API error should mark the job as failed."""
        import deepl as deepl_lib
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_translator.translate_text.side_effect = deepl_lib.DeepLException("quota exceeded")

        docx = SimpleUploadedFile("test.docx", _make_docx_bytes())

        # Run the background thread synchronously so we can inspect the result
        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        data = response.json()
        job = DocumentTranslationJob.objects.get(pk=data["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.FAILED)
        self.assertTrue(job.error_message)
        # Should NOT create a translation record on error
        self.assertEqual(Translation.objects.count(), 0)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_invalid_form_returns_json_errors(self):
        """Submitting an invalid form should return JSON with errors."""
        txt = SimpleUploadedFile("notes.txt", b"Hello", content_type="text/plain")
        response = self.client.post(reverse("document-translation"), {
            "directionChoice": "DE->EN-GB|",
            "document": txt,
            "fair_use_confirmed": "on",
        })
        self.assertEqual(response.status_code, 400)
        data = response.json()
        self.assertIn("errors", data)
        self.assertIn("document", data["errors"])

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_missing_file_returns_json_errors(self):
        """Submitting without a file should return JSON with errors."""
        response = self.client.post(reverse("document-translation"), {
            "directionChoice": "DE->EN-GB|",
        })
        self.assertEqual(response.status_code, 400)
        data = response.json()
        self.assertIn("errors", data)
        self.assertIn("document", data["errors"])

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_active_job_blocks_new_upload(self, mock_translator_cls):
        """A second upload should be rejected with 409 while a job is pending."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        docx = SimpleUploadedFile("first.docx", _make_docx_bytes())

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread") as mock_thread:
            mock_thread.return_value = MagicMock()
            resp1 = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(resp1.status_code, 200)

        # Second upload while first job is still PENDING
        docx2 = SimpleUploadedFile("second.docx", _make_docx_bytes())
        with self._patch_empty_glossary():
            resp2 = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx2,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(resp2.status_code, 409)
        data = resp2.json()
        self.assertIn("error", data)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_completed_job_allows_new_upload(self, mock_translator_cls):
        """After a job completes, a new upload should be accepted."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        docx = SimpleUploadedFile("first.docx", _make_docx_bytes())

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread") as mock_thread:
            mock_thread.return_value = MagicMock()
            resp1 = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(resp1.status_code, 200)
        # Mark first job as completed
        job = DocumentTranslationJob.objects.get(pk=resp1.json()["job_id"])
        job.status = DocumentTranslationJob.Status.COMPLETED
        job.save(update_fields=["status"])

        # Second upload should succeed now
        docx2 = SimpleUploadedFile("second.docx", _make_docx_bytes())
        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread") as mock_thread:
            mock_thread.return_value = MagicMock()
            resp2 = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx2,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(resp2.status_code, 200)


# ============================================================================
# Document Translation Helper Tests
# ============================================================================

class TranslateDocxHelperTest(TestCase):
    """Test _translate_docx helper with real docx files."""

    @patch("deeplFrontend.views.deepl.Translator")
    def test_translates_paragraph_text(self, mock_translator_cls):
        """Paragraph text in the docx should be replaced with translations."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "test.docx")
            doc = DocxDocument()
            doc.add_paragraph("Hallo Welt")
            doc.save(path)

            char_count, api_calls, _ = _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )

            self.assertGreater(char_count, 0)
            self.assertGreater(api_calls, 0)
            # Verify the file was modified
            result_doc = DocxDocument(path)
            self.assertEqual(result_doc.paragraphs[0].text, "Hello World")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_translates_table_cells(self, mock_translator_cls):
        """Table cell text should be translated."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "table.docx")
            doc = DocxDocument()
            table = doc.add_table(rows=1, cols=2)
            table.rows[0].cells[0].text = "Zelle A"
            table.rows[0].cells[1].text = "Zelle B"
            doc.save(path)

            _translate_docx(path, mock_translator, "DE", "EN-GB", None, None, None)

            result_doc = DocxDocument(path)
            self.assertEqual(result_doc.tables[0].rows[0].cells[0].text, "Translated")
            self.assertEqual(result_doc.tables[0].rows[0].cells[1].text, "Translated")

    def test_empty_paragraphs_skipped(self):
        """Empty paragraphs should not trigger API calls."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        mock_translator = MagicMock()

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "empty.docx")
            doc = DocxDocument()
            doc.add_paragraph("")  # empty
            doc.add_paragraph("   ")  # whitespace only
            doc.save(path)

            char_count, api_calls, _ = _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )
            self.assertEqual(char_count, 0)
            self.assertEqual(api_calls, 0)
            mock_translator.translate_text.assert_not_called()

    @patch("deeplFrontend.views.deepl.Translator")
    def test_multi_run_paragraph_merged(self, mock_translator_cls):
        """Multiple runs in a paragraph should be merged and translated as one."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Hello World from DeepL"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "multi_run.docx")
            doc = DocxDocument()
            para = doc.add_paragraph()
            # Simulate Word splitting text into multiple runs (editing history)
            run1 = para.add_run("Hallo ")
            run2 = para.add_run("Welt ")
            run3 = para.add_run("von DeepL")
            doc.save(path)

            _translate_docx(path, mock_translator, "DE", "EN-GB", None, None, None)

            # API should be called exactly once with the full paragraph text
            mock_translator.translate_text.assert_called_once()
            call_args = mock_translator.translate_text.call_args
            self.assertEqual(call_args[0][0], "Hallo Welt von DeepL")

            # The translated text should be in the first run, others empty
            result_doc = DocxDocument(path)
            result_runs = result_doc.paragraphs[0].runs
            self.assertEqual(result_runs[0].text, "Hello World from DeepL")
            self.assertEqual(result_runs[1].text, "")
            self.assertEqual(result_runs[2].text, "")
            # Full paragraph text should still be correct
            self.assertEqual(result_doc.paragraphs[0].text, "Hello World from DeepL")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_multi_run_spaces_preserved(self, mock_translator_cls):
        """Spaces at run boundaries should not be lost."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Hello beautiful World"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "spaces.docx")
            doc = DocxDocument()
            para = doc.add_paragraph()
            # Space belongs to first run (common pattern when formatting changes)
            para.add_run("Hallo ")
            run_bold = para.add_run("schöne ")
            run_bold.bold = True
            para.add_run("Welt")
            doc.save(path)

            _translate_docx(path, mock_translator, "DE", "EN-GB", None, None, None)

            result_doc = DocxDocument(path)
            # Full text must be intact, no missing spaces
            self.assertEqual(result_doc.paragraphs[0].text, "Hello beautiful World")

    @patch("deeplFrontend.views.deepl.Translator")
    def test_char_count_reflects_full_paragraph(self, mock_translator_cls):
        """Character count should include all run texts in the paragraph."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "count.docx")
            doc = DocxDocument()
            para = doc.add_paragraph()
            para.add_run("ABC")   # 3 chars
            para.add_run(" DE")   # 3 chars
            doc.save(path)

            char_count, api_calls, _ = _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )
            self.assertEqual(char_count, 6)
            self.assertEqual(api_calls, 1)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_linked_headers_not_translated_twice(self, mock_translator_cls):
        """Linked section headers/footers must not be translated multiple times."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "linked.docx")
            doc = DocxDocument()
            doc.add_paragraph("Body")
            s = doc.sections[0]
            s.header.is_linked_to_previous = False
            s.header.paragraphs[0].add_run("Header")
            s.footer.is_linked_to_previous = False
            s.footer.paragraphs[0].add_run("Footer")
            # Add a second section with linked (default) header/footer
            doc.add_section()
            doc.save(path)

            char_count, api_calls, _ = _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )
            # Body (4) + Header (6) + Footer (6) = 16 chars, 3 API calls
            # The linked second section should NOT add extra calls
            self.assertEqual(api_calls, 3)
            self.assertEqual(char_count, 16)

    # -----------------------------------------------------------------------
    # Helpers for footnote / endnote tests
    # -----------------------------------------------------------------------

    @staticmethod
    def _make_docx_with_notes_bytes(body_text, note_text, note_type="footnote"):
        """Return bytes of a .docx containing one user footnote or endnote.

        The note content is built by directly injecting the appropriate
        word/footnotes.xml (or word/endnotes.xml) part into the ZIP, because
        python-docx 1.1.x has no high-level API for creating footnotes.
        """
        import io
        import zipfile as zipfile_mod
        from docx import Document as DocxDocument

        buf = io.BytesIO()
        doc = DocxDocument()
        doc.add_paragraph(body_text)
        doc.save(buf)
        buf.seek(0)

        rel_type = (
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
            + ("footnotes" if note_type == "footnote" else "endnotes")
        )
        xml_file = "word/" + ("footnotes" if note_type == "footnote" else "endnotes") + ".xml"
        root_tag = "w:footnotes" if note_type == "footnote" else "w:endnotes"
        note_tag = "w:footnote" if note_type == "footnote" else "w:endnote"
        ref_tag = "w:footnoteRef" if note_type == "footnote" else "w:endnoteRef"

        notes_xml = (
            "<?xml version='1.0' encoding='UTF-8' standalone='yes'?>"
            "<{root} xmlns:w=\"http://schemas.openxmlformats.org/wordprocessingml/2006/main\">"
            "<{note} w:type=\"separator\" w:id=\"-1\">"
            "<w:p><w:r><w:separator/></w:r></w:p>"
            "</{note}>"
            "<{note} w:type=\"continuationSeparator\" w:id=\"0\">"
            "<w:p><w:r><w:continuationSeparator/></w:r></w:p>"
            "</{note}>"
            "<{note} w:id=\"1\">"
            "<w:p>"
            "<w:r><w:rPr><w:rStyle w:val=\"FootnoteReference\"/></w:rPr><{ref}/></w:r>"
            "<w:r><w:t>{text}</w:t></w:r>"
            "</w:p>"
            "</{note}>"
            "</{root}>"
        ).format(
            root=root_tag, note=note_tag, ref=ref_tag, text=note_text
        ).encode("utf-8")

        result = io.BytesIO()
        with zipfile_mod.ZipFile(buf, "r") as src, \
                zipfile_mod.ZipFile(result, "w", zipfile_mod.ZIP_DEFLATED) as dst:
            for item in src.infolist():
                data = src.read(item.filename)
                if item.filename == "word/_rels/document.xml.rels":
                    rel_entry = (
                        '<Relationship Id="rId99"'
                        ' Type="{rel_type}"'
                        ' Target="{target}"/>'
                    ).format(
                        rel_type=rel_type,
                        target=("footnotes.xml" if note_type == "footnote"
                                else "endnotes.xml"),
                    ).encode("utf-8")
                    data = data.replace(b"</Relationships>",
                                        rel_entry + b"</Relationships>")
                dst.writestr(item, data)
            dst.writestr(xml_file, notes_xml)

        result.seek(0)
        return result.getvalue()

    # -----------------------------------------------------------------------
    # P1: footnote / endnote reference mark preservation
    # -----------------------------------------------------------------------

    @patch("deeplFrontend.views.deepl.Translator")
    def test_footnote_reference_mark_preserved(self, mock_translator_cls):
        """A <w:footnoteReference> run must survive translation (P1 fix).

        python-docx's run.text setter calls CT_R.clear_content() which removes
        ALL non-rPr child elements.  Runs that carry no visible text (such as
        footnote reference mark runs) must therefore be skipped when clearing
        the multi-run paragraph after translation.
        """
        from docx import Document as DocxDocument
        from docx.oxml.parser import OxmlElement
        from docx.oxml.ns import qn
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Translated text"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "fn_ref.docx")
            doc = DocxDocument()
            para = doc.add_paragraph()
            para.add_run("Text before")
            # Append a bare footnote-reference run (no <w:t>) directly to the
            # paragraph XML — python-docx has no high-level API for this.
            fn_run = OxmlElement("w:r")
            fn_rpr = OxmlElement("w:rPr")
            fn_rstyle = OxmlElement("w:rStyle")
            fn_rstyle.set(qn("w:val"), "FootnoteReference")
            fn_rpr.append(fn_rstyle)
            fn_run.append(fn_rpr)
            fn_ref = OxmlElement("w:footnoteReference")
            fn_ref.set(qn("w:id"), "1")
            fn_run.append(fn_ref)
            para._p.append(fn_run)
            para.add_run(" text after")
            doc.save(path)

            char_count, api_calls, has_notes = _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )

            result_doc = DocxDocument(path)
            result_para = result_doc.paragraphs[0]

            # The paragraph must contain translated text.
            self.assertIn("Translated text", result_para.text)

            # The <w:footnoteReference> element must still be present.
            fn_refs = result_para._p.findall(".//" + qn("w:footnoteReference"))
            self.assertEqual(
                len(fn_refs), 1,
                "<w:footnoteReference> was destroyed during translation",
            )
            self.assertEqual(fn_refs[0].get(qn("w:id")), "1")

            # No footnotes OPC part → has_notes must be False.
            self.assertFalse(has_notes)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_endnote_reference_mark_preserved(self, mock_translator_cls):
        """A <w:endnoteReference> run must survive translation (P1 fix)."""
        from docx import Document as DocxDocument
        from docx.oxml.parser import OxmlElement
        from docx.oxml.ns import qn
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "en_ref.docx")
            doc = DocxDocument()
            para = doc.add_paragraph()
            para.add_run("Endnote text")
            en_run = OxmlElement("w:r")
            en_ref = OxmlElement("w:endnoteReference")
            en_ref.set(qn("w:id"), "1")
            en_run.append(en_ref)
            para._p.append(en_run)
            doc.save(path)

            _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )

            result_doc = DocxDocument(path)
            result_para = result_doc.paragraphs[0]
            en_refs = result_para._p.findall(".//" + qn("w:endnoteReference"))
            self.assertEqual(
                len(en_refs), 1,
                "<w:endnoteReference> was destroyed during translation",
            )

    def test_footnote_ref_run_not_counted_in_chars(self):
        """A footnote-reference run contributes 0 characters to the count."""
        from docx import Document as DocxDocument
        from docx.oxml.parser import OxmlElement
        from docx.oxml.ns import qn
        from .views import _translate_docx

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "ref_count.docx")
            doc = DocxDocument()
            para = doc.add_paragraph()
            para.add_run("ABC")      # 3 chars
            fn_run = OxmlElement("w:r")
            fn_ref = OxmlElement("w:footnoteReference")
            fn_ref.set(qn("w:id"), "1")
            fn_run.append(fn_ref)
            para._p.append(fn_run)  # 0 chars
            doc.save(path)

            char_count, api_calls, _ = _translate_docx(
                path, MagicMock(), "DE", "EN-GB", None, None, None,
                count_only=True,
            )

        self.assertEqual(char_count, 3)

    # -----------------------------------------------------------------------
    # P2: footnote / endnote content translation
    # -----------------------------------------------------------------------

    @patch("deeplFrontend.views.deepl.Translator")
    def test_footnote_content_translated(self, mock_translator_cls):
        """Footnote text content must be translated and saved back (P2 fix)."""
        import zipfile as zipfile_mod
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        docx_bytes = self._make_docx_with_notes_bytes(
            "Haupttext", "Fußnotentext", note_type="footnote",
        )
        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "footnote.docx")
            with open(path, "wb") as f:
                f.write(docx_bytes)

            char_count, api_calls, has_notes = _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )

            # has_notes must be True: user footnote was found.
            self.assertTrue(has_notes)

            # The API must have been called with the footnote text.
            all_texts = [
                call.args[0]
                for call in mock_translator.translate_text.call_args_list
            ]
            self.assertIn("Fußnotentext", all_texts)

            # The saved file must contain "Translated" inside footnotes.xml.
            with zipfile_mod.ZipFile(path, "r") as z:
                self.assertIn("word/footnotes.xml", z.namelist())
                fn_xml = z.read("word/footnotes.xml").decode("utf-8")
            self.assertIn("Translated", fn_xml)
            self.assertNotIn("Fußnotentext", fn_xml)

    @patch("deeplFrontend.views.deepl.Translator")
    def test_endnote_content_translated(self, mock_translator_cls):
        """Endnote text content must be translated and saved back (P2 fix)."""
        import zipfile as zipfile_mod
        from .views import _translate_docx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        docx_bytes = self._make_docx_with_notes_bytes(
            "Haupttext", "Endnotentext", note_type="endnote",
        )
        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "endnote.docx")
            with open(path, "wb") as f:
                f.write(docx_bytes)

            char_count, api_calls, has_notes = _translate_docx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )

            self.assertTrue(has_notes)
            all_texts = [
                call.args[0]
                for call in mock_translator.translate_text.call_args_list
            ]
            self.assertIn("Endnotentext", all_texts)

            with zipfile_mod.ZipFile(path, "r") as z:
                self.assertIn("word/endnotes.xml", z.namelist())
                en_xml = z.read("word/endnotes.xml").decode("utf-8")
            self.assertIn("Translated", en_xml)
            self.assertNotIn("Endnotentext", en_xml)

    def test_count_only_includes_footnote_chars(self):
        """count_only mode must count footnote text characters (P2 fix)."""
        from .views import _count_chars_docx

        # Body: "Hallo" (5), footnote: "Fußnotentext" (12) → total 17
        docx_bytes = self._make_docx_with_notes_bytes(
            "Hallo", "Fußnotentext", note_type="footnote",
        )
        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "count_fn.docx")
            with open(path, "wb") as f:
                f.write(docx_bytes)

            total = _count_chars_docx(path)

        # Body paragraph "Hallo" = 5 chars; footnote "Fußnotentext" = 12 chars.
        self.assertEqual(total, 17)

    def test_has_notes_false_for_plain_docx(self):
        """has_notes must be False when the document has no footnotes/endnotes."""
        from docx import Document as DocxDocument
        from .views import _translate_docx

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "plain.docx")
            doc = DocxDocument()
            doc.add_paragraph("Some text")
            doc.save(path)

            _, _, has_notes = _translate_docx(
                path, MagicMock(), "DE", "EN-GB", None, None, None,
                count_only=True,
            )

        self.assertFalse(has_notes)


class TranslatePptxHelperTest(TestCase):
    """Test _translate_pptx helper with real pptx files."""

    @patch("deeplFrontend.views.deepl.Translator")
    def test_translates_slide_text(self, mock_translator_cls):
        """Slide text frames should be translated."""
        from pptx import Presentation
        from pptx.util import Inches
        from .views import _translate_pptx

        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "test.pptx")
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Titel"
            slide.placeholders[1].text = "Hallo Welt"
            prs.save(path)

            char_count, api_calls = _translate_pptx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )

            self.assertGreater(char_count, 0)
            self.assertGreater(api_calls, 0)
            mock_translator.translate_text.assert_called()

    def test_empty_slides_skipped(self):
        """Slides with no text should not trigger API calls."""
        from pptx import Presentation
        from .views import _translate_pptx

        mock_translator = MagicMock()

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "blank.pptx")
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            prs.save(path)

            char_count, api_calls = _translate_pptx(
                path, mock_translator, "DE", "EN-GB", None, None, None,
            )

            self.assertEqual(char_count, 0)
            self.assertEqual(api_calls, 0)
            mock_translator.translate_text.assert_not_called()


class TranslateTextFragmentTest(TestCase):
    """Test the _translate_text_fragment helper."""

    def test_empty_string_returned_as_is(self):
        """Empty strings should be returned without calling the API."""
        from .views import _translate_text_fragment
        mock_translator = MagicMock()
        result = _translate_text_fragment(mock_translator, "", "DE", "EN-GB", None, None, None)
        self.assertEqual(result, "")
        mock_translator.translate_text.assert_not_called()

    def test_whitespace_only_returned_as_is(self):
        """Whitespace-only strings should be returned without calling the API."""
        from .views import _translate_text_fragment
        mock_translator = MagicMock()
        result = _translate_text_fragment(mock_translator, "   ", "DE", "EN-GB", None, None, None)
        self.assertEqual(result, "   ")
        mock_translator.translate_text.assert_not_called()

    def test_none_returned_as_is(self):
        """None should be returned without calling the API."""
        from .views import _translate_text_fragment
        mock_translator = MagicMock()
        result = _translate_text_fragment(mock_translator, None, "DE", "EN-GB", None, None, None)
        self.assertIsNone(result)
        mock_translator.translate_text.assert_not_called()

    def test_text_translated_via_api(self):
        """Non-empty text should be translated via the API."""
        from .views import _translate_text_fragment
        mock_translator = MagicMock()
        mock_result = MagicMock()
        mock_result.text = "Hello"
        mock_translator.translate_text.return_value = mock_result

        result = _translate_text_fragment(mock_translator, "Hallo", "DE", "EN-GB", None, None, None)
        self.assertEqual(result, "Hello")
        mock_translator.translate_text.assert_called_once()


class LangSuffixTest(TestCase):
    """Test the _lang_suffix helper."""

    def test_en_gb(self):
        from .views import _lang_suffix
        self.assertEqual(_lang_suffix("EN-GB"), "en")

    def test_de(self):
        from .views import _lang_suffix
        self.assertEqual(_lang_suffix("DE"), "de")

    def test_fr(self):
        from .views import _lang_suffix
        self.assertEqual(_lang_suffix("FR"), "fr")


# ============================================================================
# Document Translation Settings Test
# ============================================================================

class DocumentTranslationSettingsTest(TestCase):
    """Test DOCUMENT_TRANSLATION_ENABLED setting and context processor."""

    def test_setting_exists(self):
        """DOCUMENT_TRANSLATION_ENABLED should exist in settings."""
        self.assertTrue(hasattr(settings, "DOCUMENT_TRANSLATION_ENABLED"))

    def test_setting_is_boolean_or_string(self):
        """DOCUMENT_TRANSLATION_ENABLED should be a bool or a regex string."""
        self.assertIsInstance(settings.DOCUMENT_TRANSLATION_ENABLED, (bool, str))

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_context_processor_exposes_setting_true(self):
        """Context processor should expose DOCUMENT_TRANSLATION_ENABLED=True."""
        response = self.client.get(reverse("translation-form"))
        self.assertTrue(response.context["DOCUMENT_TRANSLATION_ENABLED"])

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_context_processor_exposes_setting_false(self):
        """Context processor should expose DOCUMENT_TRANSLATION_ENABLED=False."""
        response = self.client.get(reverse("translation-form"))
        self.assertFalse(response.context["DOCUMENT_TRANSLATION_ENABLED"])


# ============================================================================
# DocumentTranslationJob Model Tests
# ============================================================================

class DocumentTranslationJobModelTest(TestCase):
    """Test DocumentTranslationJob model."""

    def _create_job(self, **kwargs):
        defaults = {
            "session_key": "test-session-key",
            "original_filename": "test.docx",
            "file_type": ".docx",
            "file_size": 1024,
            "direction": "DE->EN-GB|",
            "output_filename": "test_en.docx",
        }
        defaults.update(kwargs)
        return DocumentTranslationJob.objects.create(**defaults)

    def test_default_status_pending(self):
        """New jobs should have status PENDING."""
        job = self._create_job()
        self.assertEqual(job.status, DocumentTranslationJob.Status.PENDING)

    def test_str_representation(self):
        """String representation includes filename and status."""
        job = self._create_job()
        self.assertIn("test.docx", str(job))
        self.assertIn("pending", str(job))

    def test_uuid_primary_key(self):
        """Primary key should be a valid UUID."""
        import uuid as uuid_mod
        job = self._create_job()
        self.assertIsInstance(job.pk, uuid_mod.UUID)

    def test_downloaded_default_false(self):
        """downloaded should default to False."""
        job = self._create_job()
        self.assertFalse(job.downloaded)


# ============================================================================
# Document Job Status Endpoint Tests
# ============================================================================

class DocumentJobStatusEndpointTest(TestCase):
    """Test the /translation/document/<id>/status/ polling endpoint."""

    def setUp(self):
        self.client = Client()
        # Force a session so we have a session_key
        self.client.get(reverse("translation-form"))
        session = self.client.session
        session.save()
        self.session_key = session.session_key

    def _create_job(self, **kwargs):
        defaults = {
            "session_key": self.session_key,
            "original_filename": "test.docx",
            "file_type": ".docx",
            "file_size": 2048,
            "direction": "DE->EN-GB|",
            "output_filename": "test_en.docx",
        }
        defaults.update(kwargs)
        return DocumentTranslationJob.objects.create(**defaults)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_pending_status(self):
        """Pending job should return status 'pending'."""
        job = self._create_job()
        response = self.client.get(reverse("document-job-status", args=[job.id]))
        self.assertEqual(response.status_code, 200)
        data = response.json()
        self.assertEqual(data["status"], "pending")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_completed_status_includes_metadata(self):
        """Completed job should include characters, api_calls, duration."""
        job = self._create_job(
            status=DocumentTranslationJob.Status.COMPLETED,
            characters=500,
            api_calls=10,
            duration_seconds=2.5,
        )
        response = self.client.get(reverse("document-job-status", args=[job.id]))
        data = response.json()
        self.assertEqual(data["status"], "completed")
        self.assertEqual(data["characters"], 500)
        self.assertEqual(data["api_calls"], 10)
        self.assertEqual(data["duration_seconds"], 2.5)
        self.assertIn("has_footnotes", data)
        self.assertFalse(data["has_footnotes"])

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_completed_status_has_footnotes_true(self):
        """Status endpoint must report has_footnotes=True when set on the job."""
        job = self._create_job(
            status=DocumentTranslationJob.Status.COMPLETED,
            characters=100,
            api_calls=2,
            duration_seconds=1.0,
            has_footnotes=True,
        )
        response = self.client.get(reverse("document-job-status", args=[job.id]))
        data = response.json()
        self.assertEqual(data["status"], "completed")
        self.assertTrue(data["has_footnotes"])

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_failed_status_includes_error(self):
        """Failed job should include error_message."""
        job = self._create_job(
            status=DocumentTranslationJob.Status.FAILED,
            error_message="API quota exceeded",
        )
        response = self.client.get(reverse("document-job-status", args=[job.id]))
        data = response.json()
        self.assertEqual(data["status"], "failed")
        self.assertEqual(data["error_message"], "API quota exceeded")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_wrong_session_returns_403(self):
        """Polling a job from another session should return 403."""
        job = self._create_job(session_key="other-session-key")
        response = self.client.get(reverse("document-job-status", args=[job.id]))
        self.assertEqual(response.status_code, 403)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_nonexistent_job_returns_404(self):
        """Polling a non-existent job should return 404."""
        import uuid as uuid_mod
        fake_id = uuid_mod.uuid4()
        response = self.client.get(reverse("document-job-status", args=[fake_id]))
        self.assertEqual(response.status_code, 404)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_disabled_returns_404(self):
        """Status endpoint should return 404 when feature is disabled."""
        job = self._create_job()
        response = self.client.get(reverse("document-job-status", args=[job.id]))
        self.assertEqual(response.status_code, 404)


# ============================================================================
# Document Download Endpoint Tests
# ============================================================================

class DocumentDownloadEndpointTest(TestCase):
    """Test the /translation/document/<id>/download/ endpoint."""

    def setUp(self):
        self.client = Client()
        self.client.get(reverse("translation-form"))
        session = self.client.session
        session.save()
        self.session_key = session.session_key
        self.tmp_dir = tempfile.mkdtemp()

    def tearDown(self):
        import shutil
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def _create_completed_job(self, with_file=True, **kwargs):
        defaults = {
            "session_key": self.session_key,
            "original_filename": "report.docx",
            "file_type": ".docx",
            "file_size": 4096,
            "direction": "DE->EN-GB|",
            "output_filename": "report_en.docx",
            "status": DocumentTranslationJob.Status.COMPLETED,
            "characters": 100,
        }
        defaults.update(kwargs)
        job = DocumentTranslationJob.objects.create(**defaults)

        if with_file:
            job_dir = os.path.join(self.tmp_dir, "doc_translations", str(job.id))
            os.makedirs(job_dir, exist_ok=True)
            filepath = os.path.join(job_dir, "report.docx")
            with open(filepath, "wb") as f:
                f.write(b"PK\x03\x04fake-docx-content")
            job.result_path = os.path.join("doc_translations", str(job.id), "report.docx")
            job.save(update_fields=["result_path"])

        return job

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_download_serves_file(self):
        """Download should serve the file with correct Content-Disposition."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job()
            response = self.client.get(reverse("document-download", args=[job.id]))
        self.assertEqual(response.status_code, 200)
        self.assertIn("report_en.docx", response["Content-Disposition"])

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_download_marks_as_downloaded(self):
        """Download should set downloaded=True and increment download_count."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job()
            self.client.get(reverse("document-download", args=[job.id]))
        job.refresh_from_db()
        self.assertTrue(job.downloaded)
        self.assertEqual(job.download_count, 1)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_download_count_increments(self):
        """Each download should increment download_count."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job()
            self.client.get(reverse("document-download", args=[job.id]))
            self.client.get(reverse("document-download", args=[job.id]))
        job.refresh_from_db()
        self.assertEqual(job.download_count, 2)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_download_keeps_file_for_redownload(self):
        """File should remain on disk after download (10-min cleanup handles removal)."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job()
            self.client.get(reverse("document-download", args=[job.id]))
            full_path = os.path.join(self.tmp_dir, job.result_path)
            self.assertTrue(os.path.isfile(full_path))

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_wrong_session_returns_403(self):
        """Downloading a job from another session should return 403."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job(session_key="other-session")
            response = self.client.get(reverse("document-download", args=[job.id]))
        self.assertEqual(response.status_code, 403)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_pending_job_returns_409(self):
        """Downloading a non-completed job should return 409."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job(
                status=DocumentTranslationJob.Status.PENDING,
            )
            response = self.client.get(reverse("document-download", args=[job.id]))
        self.assertEqual(response.status_code, 409)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_missing_result_path_returns_410(self):
        """Downloading a job with empty result_path should return 410."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job(with_file=False)
            response = self.client.get(reverse("document-download", args=[job.id]))
        self.assertEqual(response.status_code, 410)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_disabled_returns_404(self):
        """Download endpoint should return 404 when feature is disabled."""
        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_completed_job()
            response = self.client.get(reverse("document-download", args=[job.id]))
        self.assertEqual(response.status_code, 404)


# ============================================================================
# Stale Document Job Cleanup Tests
# ============================================================================

class StaleDocumentJobCleanupTest(TestCase):
    """Test _cleanup_stale_document_jobs helper."""

    def setUp(self):
        self.tmp_dir = tempfile.mkdtemp()

    def tearDown(self):
        import shutil
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def _create_job_with_file(self, created_minutes_ago=15, **kwargs):
        defaults = {
            "session_key": "cleanup-test",
            "original_filename": "old.docx",
            "file_type": ".docx",
            "file_size": 512,
            "direction": "DE->EN-GB|",
            "output_filename": "old_en.docx",
            "status": DocumentTranslationJob.Status.COMPLETED,
        }
        defaults.update(kwargs)
        job = DocumentTranslationJob.objects.create(**defaults)

        # Backdate created_at
        DocumentTranslationJob.objects.filter(pk=job.pk).update(
            created_at=timezone.now() - timedelta(minutes=created_minutes_ago),
        )
        job.refresh_from_db()

        # Create a file on disk
        job_dir = os.path.join(self.tmp_dir, "doc_translations", str(job.id))
        os.makedirs(job_dir, exist_ok=True)
        filepath = os.path.join(job_dir, "old.docx")
        with open(filepath, "wb") as f:
            f.write(b"fake content")
        job.result_path = os.path.join("doc_translations", str(job.id), "old.docx")
        job.save(update_fields=["result_path"])
        return job

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_stale_jobs_cleaned(self):
        """Jobs older than 10 minutes should have their files removed."""
        from .views import _cleanup_stale_document_jobs

        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_job_with_file(created_minutes_ago=15)
            filepath = os.path.join(self.tmp_dir, job.result_path)
            self.assertTrue(os.path.isfile(filepath))

            _cleanup_stale_document_jobs()

            self.assertFalse(os.path.isfile(filepath))
            job.refresh_from_db()
            self.assertEqual(job.result_path, "")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_recent_jobs_kept(self):
        """Jobs younger than 10 minutes should not be cleaned up."""
        from .views import _cleanup_stale_document_jobs

        with self.settings(MEDIA_ROOT=self.tmp_dir):
            job = self._create_job_with_file(created_minutes_ago=5)
            filepath = os.path.join(self.tmp_dir, job.result_path)

            _cleanup_stale_document_jobs()

            self.assertTrue(os.path.isfile(filepath))
            job.refresh_from_db()
            self.assertNotEqual(job.result_path, "")


# ============================================================================
# Document Translation Timeout Tests
# ============================================================================

class DocumentTranslationTimeoutTest(TestCase):
    """Test the DOCUMENT_TRANSLATION_TIMEOUT wall-time limit."""

    def setUp(self):
        self.client = Client()
        from .views import _GlossaryCache
        self._empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        self._empty_cache._cache = {}
        self._empty_cache._loaded_at = float("inf")

    def _patch_empty_glossary(self):
        return patch("deeplFrontend.views.glossary_cache", self._empty_cache)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True, DOCUMENT_TRANSLATION_TIMEOUT=1)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_timeout_marks_job_failed(self, mock_translator_cls):
        """A translation exceeding the wall-time limit should fail."""
        import time as time_mod

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        # Simulate a slow API call that takes longer than the 1-second limit
        def slow_translate(*args, **kwargs):
            time_mod.sleep(0.3)
            result = MagicMock()
            result.text = "Translated"
            return result
        mock_translator.translate_text.side_effect = slow_translate

        # Build a docx with enough paragraphs to force multiple API calls
        from docx import Document as DocxDocument
        buf = BytesIO()
        doc = DocxDocument()
        for i in range(10):
            doc.add_paragraph(f"Paragraph {i} with some text to translate for testing.")
        doc.save(buf)
        docx_bytes = buf.getvalue()

        docx = SimpleUploadedFile("slow.docx", docx_bytes)

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.FAILED)
        self.assertIn("1", job.error_message)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True, DOCUMENT_TRANSLATION_TIMEOUT=0)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_timeout_disabled_allows_completion(self, mock_translator_cls):
        """With timeout=0, translation should complete normally."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Translated"
        mock_translator.translate_text.return_value = mock_result

        docx = SimpleUploadedFile("fast.docx", _make_docx_bytes())

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.COMPLETED)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_default_timeout_setting(self):
        """DOCUMENT_TRANSLATION_TIMEOUT should default to 180."""
        self.assertEqual(settings.DOCUMENT_TRANSLATION_TIMEOUT, 180)


# ============================================================================
# Document Character Limit Tests
# ============================================================================

class DocumentCharacterLimitTest(TestCase):
    """Test pre-flight character count check for document translations."""

    def setUp(self):
        self.client = Client()
        from .views import _GlossaryCache
        self._empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        self._empty_cache._cache = {}
        self._empty_cache._loaded_at = float("inf")

    def _patch_empty_glossary(self):
        return patch("deeplFrontend.views.glossary_cache", self._empty_cache)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True, MAX_TRANSLATION_LENGTH=9)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_exceeding_char_limit_fails_job(self, mock_translator_cls):
        """A document exceeding MAX_TRANSLATION_LENGTH should fail immediately."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        docx = SimpleUploadedFile("big.docx", _make_docx_bytes())  # "Hallo Welt" = 10 chars

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.FAILED)
        self.assertIn("9", job.error_message)
        # Should NOT have called the translation API
        mock_translator.translate_text.assert_not_called()

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True, MAX_TRANSLATION_LENGTH=50000)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_within_char_limit_succeeds(self, mock_translator_cls):
        """A document within MAX_TRANSLATION_LENGTH should translate normally."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        docx = SimpleUploadedFile("small.docx", _make_docx_bytes())

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.COMPLETED)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True, MAX_TRANSLATION_LENGTH=0)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_limit_disabled_allows_any_size(self, mock_translator_cls):
        """With MAX_TRANSLATION_LENGTH=0, any document should be accepted."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_translator.translate_text.return_value = mock_result

        docx = SimpleUploadedFile("any.docx", _make_docx_bytes())

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.COMPLETED)

    def test_count_chars_docx(self):
        """_count_chars_docx should count characters in a .docx file."""
        from docx import Document as DocxDocument
        from .views import _count_chars_docx

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "test.docx")
            doc = DocxDocument()
            doc.add_paragraph("Hallo Welt")  # 10 chars
            doc.add_paragraph("Noch mehr Text")  # 14 chars
            doc.save(path)

            count = _count_chars_docx(path)
            self.assertEqual(count, 24)

    def test_count_chars_docx_empty(self):
        """Empty docx should return 0 characters."""
        from docx import Document as DocxDocument
        from .views import _count_chars_docx

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "empty.docx")
            doc = DocxDocument()
            doc.add_paragraph("")
            doc.save(path)

            count = _count_chars_docx(path)
            self.assertEqual(count, 0)

    def test_count_chars_pptx(self):
        """_count_chars_pptx should count characters in a .pptx file."""
        from pptx import Presentation
        from pptx.util import Inches
        from .views import _count_chars_pptx

        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "test.pptx")
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Titel"  # 5 chars
            slide.placeholders[1].text = "Hallo Welt"  # 10 chars
            prs.save(path)

            count = _count_chars_pptx(path)
            self.assertEqual(count, 15)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True, MAX_TRANSLATION_LENGTH=10)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_error_message_contains_exceeded_percentage(self, mock_translator_cls):
        """Error message should contain the allowed limit and exceeded percentage."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        # Build docx with lots of text (well over 10 chars)
        from docx import Document as DocxDocument
        buf = BytesIO()
        doc = DocxDocument()
        doc.add_paragraph("This is a much longer paragraph that exceeds ten characters definitely")
        doc.save(buf)
        docx = SimpleUploadedFile("over.docx", buf.getvalue())

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.FAILED)
        # Error message should mention the limit (10) and a percentage
        self.assertIn("10", job.error_message)
        self.assertIn("%", job.error_message)


# ============================================================================
# File Cleanup on Failure Tests
# ============================================================================

class FailedJobFileCleanupTest(TestCase):
    """Test that files are deleted immediately when a job fails."""

    def setUp(self):
        self.client = Client()
        from .views import _GlossaryCache
        self._empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        self._empty_cache._cache = {}
        self._empty_cache._loaded_at = float("inf")

    def _patch_empty_glossary(self):
        return patch("deeplFrontend.views.glossary_cache", self._empty_cache)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_api_error_deletes_file(self, mock_translator_cls):
        """A DeepL API error should delete the uploaded file from disk."""
        import deepl as deepl_lib
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator
        mock_translator.translate_text.side_effect = deepl_lib.DeepLException("error")

        docx = SimpleUploadedFile("secret.docx", _make_docx_bytes())

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.FAILED)
        # result_path should be cleared (file deleted)
        self.assertEqual(job.result_path, "")

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True, MAX_TRANSLATION_LENGTH=5)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_char_limit_exceeded_deletes_file(self, mock_translator_cls):
        """Exceeding char limit should delete the uploaded file from disk."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        docx = SimpleUploadedFile("big.docx", _make_docx_bytes())

        def run_sync(**kwargs):
            target = kwargs["target"]
            t_args = kwargs.get("args", ())
            mock_t = MagicMock()
            mock_t.start = lambda: target(*t_args)
            return mock_t

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread", side_effect=run_sync):
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        self.assertEqual(job.status, DocumentTranslationJob.Status.FAILED)
        self.assertEqual(job.result_path, "")
        mock_translator.translate_text.assert_not_called()


# ============================================================================
# Stale Failed Job Cleanup Test
# ============================================================================

class StaleFailedJobCleanupTest(TestCase):
    """Test that failed job files are also cleaned up by the lazy cleanup."""

    def setUp(self):
        self.tmp_dir = tempfile.mkdtemp()

    def tearDown(self):
        import shutil
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_stale_failed_job_cleaned(self):
        """A failed job older than 10 min that still has result_path should be cleaned."""
        from .views import _cleanup_stale_document_jobs

        job = DocumentTranslationJob.objects.create(
            session_key="cleanup-fail-test",
            original_filename="fail.docx",
            file_type=".docx",
            file_size=512,
            direction="DE->EN-GB|",
            output_filename="fail_en.docx",
            status=DocumentTranslationJob.Status.FAILED,
            error_message="Some error",
        )
        # Backdate
        DocumentTranslationJob.objects.filter(pk=job.pk).update(
            created_at=timezone.now() - timedelta(minutes=15),
        )
        job.refresh_from_db()

        # Create file on disk
        job_dir = os.path.join(self.tmp_dir, "doc_translations", str(job.id))
        os.makedirs(job_dir, exist_ok=True)
        filepath = os.path.join(job_dir, "fail.docx")
        with open(filepath, "wb") as f:
            f.write(b"sensitive content")
        job.result_path = os.path.join("doc_translations", str(job.id), "fail.docx")
        job.save(update_fields=["result_path"])

        with self.settings(MEDIA_ROOT=self.tmp_dir):
            _cleanup_stale_document_jobs()

        self.assertFalse(os.path.isfile(filepath))
        job.refresh_from_db()
        self.assertEqual(job.result_path, "")


# ============================================================================
# Filename Sanitisation Test
# ============================================================================

class FilenameSanitisationTest(TestCase):
    """Test that uploaded filenames are sanitised to prevent path-traversal."""

    def setUp(self):
        self.client = Client()
        from .views import _GlossaryCache
        self._empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        self._empty_cache._cache = {}
        self._empty_cache._loaded_at = float("inf")

    def _patch_empty_glossary(self):
        return patch("deeplFrontend.views.glossary_cache", self._empty_cache)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_path_traversal_stripped(self, mock_translator_cls):
        """Directory components in the filename should be stripped."""
        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        docx_bytes = _make_docx_bytes()
        docx = SimpleUploadedFile("../../etc/evil.docx", docx_bytes)

        with self._patch_empty_glossary(), \
             patch("deeplFrontend.views.threading.Thread") as mock_thread:
            mock_thread.return_value = MagicMock()
            response = self.client.post(reverse("document-translation"), {
                "directionChoice": "DE->EN-GB|",
                "document": docx,
                "fair_use_confirmed": "on",
            })
        self.assertEqual(response.status_code, 200)
        job = DocumentTranslationJob.objects.get(pk=response.json()["job_id"])
        # The job should store only the base filename
        self.assertEqual(job.original_filename, "evil.docx")
        self.assertNotIn("..", job.result_path)
        self.assertNotIn("..", job.output_filename)


# ============================================================================
# Feature-Off Isolation Test
# ============================================================================

class FeatureOffIsolationTest(TestCase):
    """Verify the app works normally with DOCUMENT_TRANSLATION_ENABLED=False."""

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    @patch("deeplFrontend.views.deepl.Translator")
    def test_text_translation_works_when_doc_feature_off(self, mock_translator_cls):
        """Text translation should function normally with doc feature disabled."""
        from .views import _GlossaryCache
        empty_cache = _GlossaryCache.__new__(_GlossaryCache)
        empty_cache._cache = {}
        empty_cache._loaded_at = float("inf")

        mock_translator = MagicMock()
        mock_translator_cls.return_value = mock_translator

        mock_result = MagicMock()
        mock_result.text = "Hello World"
        mock_result.detected_source_lang = "DE"
        mock_translator.translate_text.return_value = mock_result

        mock_usage = MagicMock()
        mock_usage.character.count = 100
        mock_usage.character.limit = 500000
        mock_translator.get_usage.return_value = mock_usage

        client = Client()
        with patch("deeplFrontend.views.glossary_cache", empty_cache):
            response = client.post(reverse("translation-form"), {
                "sourceText": "Hallo Welt",
                "directionChoice": "DE->EN-GB|",
            })

        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Hello World")
        # No doc_form in context when disabled? Actually doc_form is always passed.
        # But no active_doc_job_json
        self.assertNotIn("active_doc_job_json", response.context)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_no_cleanup_runs_when_disabled(self):
        """_cleanup_stale_document_jobs should not be called when feature is off."""
        client = Client()
        with patch("deeplFrontend.views._cleanup_stale_document_jobs") as mock_cleanup:
            client.get(reverse("translation-form"))
        mock_cleanup.assert_not_called()

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_all_doc_endpoints_return_404_or_405(self):
        """Document endpoints should be inaccessible when feature is off."""
        client = Client()
        # POST to document-translation
        docx = SimpleUploadedFile("test.docx", _make_docx_bytes())
        resp = client.post(reverse("document-translation"), {
            "directionChoice": "DE->EN-GB|",
            "document": docx,
            "fair_use_confirmed": "on",
        })
        self.assertEqual(resp.status_code, 404)

        # GET to status endpoint
        import uuid
        fake_id = str(uuid.uuid4())
        resp = client.get(reverse("document-job-status", args=[fake_id]))
        self.assertEqual(resp.status_code, 404)

        # GET to download endpoint
        resp = client.get(reverse("document-download", args=[fake_id]))
        self.assertEqual(resp.status_code, 404)


# ============================================================================
# Document Translations Middleware Tests
# ============================================================================


class DocumentTranslationsMiddlewareTest(TestCase):
    """Unit tests for DocumentTranslationsMiddleware."""

    def _get_request_with_session(self, email=None, session_data=None):
        """Return a minimal request with a dict-based session."""
        factory = RequestFactory()
        req = factory.get("/")
        req.session = dict(session_data or {})
        if email is not None:
            req.META['HTTP_X_REMOTE_EMAIL'] = email
        return req

    def _run_middleware(self, request):
        """Run the middleware and return the request (session may be mutated)."""
        from config.middleware import DocumentTranslationsMiddleware
        from django.http import HttpResponse

        def dummy_response(req):
            return HttpResponse()

        mw = DocumentTranslationsMiddleware(dummy_response)
        mw(request)
        return request

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=r'@university\.org$')
    def test_matching_email_sets_session_flag(self):
        """When the email matches the regex the session flag is set to True."""
        req = self._get_request_with_session(email='alice@university.org')
        self._run_middleware(req)
        self.assertTrue(req.session.get('document_translations'))

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=r'@university\.org$')
    def test_non_matching_email_removes_session_flag(self):
        """When the email does not match, any existing flag is cleared."""
        req = self._get_request_with_session(
            email='outsider@example.com',
            session_data={'document_translations': True},
        )
        self._run_middleware(req)
        self.assertNotIn('document_translations', req.session)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=True)
    def test_boolean_true_does_not_touch_session(self):
        """When the setting is True the middleware is a no-op."""
        req = self._get_request_with_session(email='alice@university.org')
        self._run_middleware(req)
        self.assertNotIn('document_translations', req.session)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=False)
    def test_boolean_false_does_not_touch_session(self):
        """When the setting is False the middleware is a no-op."""
        req = self._get_request_with_session(email='alice@university.org')
        self._run_middleware(req)
        self.assertNotIn('document_translations', req.session)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=r'@university\.org$')
    def test_no_email_header_leaves_session_unchanged(self):
        """Without the X-Remote-Email header the session is not modified."""
        req = self._get_request_with_session()  # no email kwarg
        self._run_middleware(req)
        self.assertNotIn('document_translations', req.session)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=r'@UNIVERSITY\.ORG$')
    def test_matching_is_case_insensitive(self):
        """Regex matching is case-insensitive."""
        req = self._get_request_with_session(email='alice@university.org')
        self._run_middleware(req)
        self.assertTrue(req.session.get('document_translations'))

    @override_settings(DOCUMENT_TRANSLATION_ENABLED='[invalid(regex')
    def test_invalid_regex_does_not_crash(self):
        """An invalid regex logs a warning and lets the request pass through."""
        req = self._get_request_with_session(email='alice@university.org')
        # Should not raise; session should be untouched.
        self._run_middleware(req)
        self.assertNotIn('document_translations', req.session)

    @override_settings(DOCUMENT_TRANSLATION_ENABLED=r'@university\.org$')
    def test_flag_already_set_stays_set(self):
        """Matching email when flag is already True does not cause issues."""
        req = self._get_request_with_session(
            email='alice@university.org',
            session_data={'document_translations': True},
        )
        self._run_middleware(req)
        self.assertTrue(req.session.get('document_translations'))


# ============================================================================
# Doc Translation Helper Function Tests
# ============================================================================


class DocTranslationHelperTest(TestCase):
    """Unit tests for _is_doc_translation_possible and _is_doc_translation_enabled."""

    def test_possible_false(self):
        """_is_doc_translation_possible returns False when setting is False."""
        from .views import _is_doc_translation_possible
        with self.settings(DOCUMENT_TRANSLATION_ENABLED=False):
            self.assertFalse(_is_doc_translation_possible())

    def test_possible_true(self):
        """_is_doc_translation_possible returns True when setting is True."""
        from .views import _is_doc_translation_possible
        with self.settings(DOCUMENT_TRANSLATION_ENABLED=True):
            self.assertTrue(_is_doc_translation_possible())

    def test_possible_regex_string(self):
        """_is_doc_translation_possible returns True when setting is a regex string."""
        from .views import _is_doc_translation_possible
        with self.settings(DOCUMENT_TRANSLATION_ENABLED=r'@university\.org$'):
            self.assertTrue(_is_doc_translation_possible())

    def test_enabled_false_always_returns_false(self):
        """_is_doc_translation_enabled returns False when setting is False."""
        from .views import _is_doc_translation_enabled
        factory = RequestFactory()
        with self.settings(DOCUMENT_TRANSLATION_ENABLED=False):
            req = factory.get("/")
            req.session = {'document_translations': True}
            self.assertFalse(_is_doc_translation_enabled(req))

    def test_enabled_true_always_returns_true(self):
        """_is_doc_translation_enabled returns True when setting is True."""
        from .views import _is_doc_translation_enabled
        factory = RequestFactory()
        with self.settings(DOCUMENT_TRANSLATION_ENABLED=True):
            req = factory.get("/")
            req.session = {}
            self.assertTrue(_is_doc_translation_enabled(req))

    def test_enabled_regex_with_session_flag(self):
        """Regex mode + session flag → True."""
        from .views import _is_doc_translation_enabled
        factory = RequestFactory()
        with self.settings(DOCUMENT_TRANSLATION_ENABLED=r'@university\.org$'):
            req = factory.get("/")
            req.session = {'document_translations': True}
            self.assertTrue(_is_doc_translation_enabled(req))

    def test_enabled_regex_without_session_flag(self):
        """Regex mode + no session flag → False."""
        from .views import _is_doc_translation_enabled
        factory = RequestFactory()
        with self.settings(DOCUMENT_TRANSLATION_ENABLED=r'@university\.org$'):
            req = factory.get("/")
            req.session = {}
            self.assertFalse(_is_doc_translation_enabled(req))

